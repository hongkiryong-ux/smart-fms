# -*- coding: utf-8 -*-
"""Smart FMS QSS형 발표자료 PPT 생성 (타 플랫폼 QSS 발표 구조 반영)."""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "docs" / "Smart_FMS_발표자료.pptx"

POSCO_BLUE = RGBColor(0x00, 0x38, 0x76)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x22, 0x22, 0x22)
GRAY = RGBColor(0x55, 0x55, 0x55)
LIGHT = RGBColor(0xF0, 0xF4, 0xF8)
ACCENT = RGBColor(0x00, 0x6E, 0xB8)
GOLD = RGBColor(0xC4, 0xA3, 0x5A)


def _bg(slide, color: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _bar(slide, prs) -> None:
    shape = slide.shapes.add_shape(1, Inches(0), Inches(0), prs.slide_width, Inches(0.42))
    shape.fill.solid()
    shape.fill.fore_color.rgb = POSCO_BLUE
    shape.line.fill.background()


def _footer(slide, text: str = "POSCO WIDE  ·  Smart FMS") -> None:
    box = slide.shapes.add_textbox(Inches(0.6), Inches(7.05), Inches(12), Inches(0.35))
    p = box.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(11)
    p.font.color.rgb = GRAY


def _title_cover(prs, title: str, lines: list[str]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, POSCO_BLUE)
    accent = slide.shapes.add_shape(1, Inches(0), Inches(6.7), prs.slide_width, Inches(0.8))
    accent.fill.solid()
    accent.fill.fore_color.rgb = ACCENT
    accent.line.fill.background()

    tag = slide.shapes.add_textbox(Inches(0.8), Inches(1.3), Inches(11), Inches(0.4))
    tag.text_frame.text = "WIDE:A Day  ·  QSS 제안형 발표자료"
    tag.text_frame.paragraphs[0].font.size = Pt(14)
    tag.text_frame.paragraphs[0].font.color.rgb = GOLD

    t = slide.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(11.5), Inches(2.2))
    tf = t.text_frame
    tf.word_wrap = True
    tf.text = title
    tf.paragraphs[0].font.size = Pt(36)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = WHITE

    sub = slide.shapes.add_textbox(Inches(0.8), Inches(4.5), Inches(11.5), Inches(1.8))
    stf = sub.text_frame
    stf.word_wrap = True
    for i, line in enumerate(lines):
        p = stf.paragraphs[0] if i == 0 else stf.add_paragraph()
        p.text = line
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(0xCC, 0xDD, 0xEE)
        p.space_after = Pt(6)


def _section(prs, num: str, title: str, subtitle: str = "") -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, LIGHT)
    _bar(slide, prs)
    n = slide.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(3), Inches(0.5))
    n.text_frame.text = num
    n.text_frame.paragraphs[0].font.size = Pt(18)
    n.text_frame.paragraphs[0].font.bold = True
    n.text_frame.paragraphs[0].font.color.rgb = ACCENT
    t = slide.shapes.add_textbox(Inches(0.8), Inches(2.6), Inches(11.5), Inches(1.5))
    t.text_frame.text = title
    t.text_frame.paragraphs[0].font.size = Pt(34)
    t.text_frame.paragraphs[0].font.bold = True
    t.text_frame.paragraphs[0].font.color.rgb = POSCO_BLUE
    if subtitle:
        s = slide.shapes.add_textbox(Inches(0.8), Inches(4.3), Inches(11), Inches(1))
        s.text_frame.text = subtitle
        s.text_frame.paragraphs[0].font.size = Pt(16)
        s.text_frame.paragraphs[0].font.color.rgb = GRAY


def _content(prs, title: str, bullets: list[str], note: str = "") -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, WHITE)
    _bar(slide, prs)
    t = slide.shapes.add_textbox(Inches(0.6), Inches(0.6), Inches(12), Inches(0.65))
    t.text_frame.text = title
    t.text_frame.paragraphs[0].font.size = Pt(26)
    t.text_frame.paragraphs[0].font.bold = True
    t.text_frame.paragraphs[0].font.color.rgb = POSCO_BLUE

    body = slide.shapes.add_textbox(Inches(0.7), Inches(1.4), Inches(12), Inches(5.3))
    tf = body.text_frame
    tf.word_wrap = True
    for i, line in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.space_after = Pt(7)
        if line.startswith("  ·") or line.startswith("   "):
            p.level = 1
            p.font.size = Pt(14)
            p.font.color.rgb = GRAY
        elif line.endswith(":") or (line and not line.startswith(" ") and len(line) < 40 and i > 0 and bullets[i - 1] == ""):
            p.font.size = Pt(16)
            p.font.bold = True
            p.font.color.rgb = ACCENT
        else:
            p.font.size = Pt(15)
            p.font.color.rgb = DARK
        if line == "":
            p.font.size = Pt(8)
    if note:
        n = slide.shapes.add_textbox(Inches(0.7), Inches(6.7), Inches(12), Inches(0.35))
        n.text_frame.text = f"※ {note}"
        n.text_frame.paragraphs[0].font.size = Pt(11)
        n.text_frame.paragraphs[0].font.color.rgb = GRAY
    else:
        _footer(slide)


def _two_col(prs, title: str, left_h: str, left: list[str], right_h: str, right: list[str]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, WHITE)
    _bar(slide, prs)
    t = slide.shapes.add_textbox(Inches(0.6), Inches(0.6), Inches(12), Inches(0.55))
    t.text_frame.text = title
    t.text_frame.paragraphs[0].font.size = Pt(24)
    t.text_frame.paragraphs[0].font.bold = True
    t.text_frame.paragraphs[0].font.color.rgb = POSCO_BLUE

    for x, h, items in ((0.6, left_h, left), (6.7, right_h, right)):
        card = slide.shapes.add_shape(1, Inches(x), Inches(1.3), Inches(5.8), Inches(5.3))
        card.fill.solid()
        card.fill.fore_color.rgb = LIGHT
        card.line.fill.background()
        ht = slide.shapes.add_textbox(Inches(x + 0.25), Inches(1.45), Inches(5.3), Inches(0.45))
        ht.text_frame.text = h
        ht.text_frame.paragraphs[0].font.size = Pt(16)
        ht.text_frame.paragraphs[0].font.bold = True
        ht.text_frame.paragraphs[0].font.color.rgb = POSCO_BLUE
        body = slide.shapes.add_textbox(Inches(x + 0.25), Inches(2.05), Inches(5.3), Inches(4.3))
        tf = body.text_frame
        tf.word_wrap = True
        for i, line in enumerate(items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = line
            p.font.size = Pt(13)
            p.font.color.rgb = DARK
            p.space_after = Pt(6)
    _footer(slide)


def _metrics(prs, title: str, cards: list[tuple[str, str, str]]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, WHITE)
    _bar(slide, prs)
    t = slide.shapes.add_textbox(Inches(0.6), Inches(0.6), Inches(12), Inches(0.55))
    t.text_frame.text = title
    t.text_frame.paragraphs[0].font.size = Pt(24)
    t.text_frame.paragraphs[0].font.bold = True
    t.text_frame.paragraphs[0].font.color.rgb = POSCO_BLUE

    n = len(cards)
    width = 12.0 / n - 0.15
    for i, (num, label, desc) in enumerate(cards):
        x = 0.55 + i * (width + 0.15)
        card = slide.shapes.add_shape(1, Inches(x), Inches(1.5), Inches(width), Inches(4.6))
        card.fill.solid()
        card.fill.fore_color.rgb = POSCO_BLUE if i % 2 == 0 else ACCENT
        card.line.fill.background()
        nb = slide.shapes.add_textbox(Inches(x + 0.15), Inches(2.1), Inches(width - 0.3), Inches(1.2))
        nb.text_frame.text = num
        nb.text_frame.paragraphs[0].font.size = Pt(36)
        nb.text_frame.paragraphs[0].font.bold = True
        nb.text_frame.paragraphs[0].font.color.rgb = WHITE
        nb.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        lb = slide.shapes.add_textbox(Inches(x + 0.15), Inches(3.4), Inches(width - 0.3), Inches(0.6))
        lb.text_frame.text = label
        lb.text_frame.paragraphs[0].font.size = Pt(16)
        lb.text_frame.paragraphs[0].font.bold = True
        lb.text_frame.paragraphs[0].font.color.rgb = WHITE
        lb.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        db = slide.shapes.add_textbox(Inches(x + 0.2), Inches(4.2), Inches(width - 0.4), Inches(1.5))
        db.text_frame.word_wrap = True
        db.text_frame.text = desc
        db.text_frame.paragraphs[0].font.size = Pt(12)
        db.text_frame.paragraphs[0].font.color.rgb = RGBColor(0xDD, 0xEE, 0xFF)
        db.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    _footer(slide)


def _table(prs, title: str, headers: list[str], rows: list[list[str]]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, WHITE)
    _bar(slide, prs)
    t = slide.shapes.add_textbox(Inches(0.6), Inches(0.55), Inches(12), Inches(0.5))
    t.text_frame.text = title
    t.text_frame.paragraphs[0].font.size = Pt(22)
    t.text_frame.paragraphs[0].font.bold = True
    t.text_frame.paragraphs[0].font.color.rgb = POSCO_BLUE

    cols = len(headers)
    table = slide.shapes.add_table(
        len(rows) + 1, cols, Inches(0.45), Inches(1.25), Inches(12.4), Inches(0.42 * (len(rows) + 1.5))
    ).table
    for ci, h in enumerate(headers):
        cell = table.cell(0, ci)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = POSCO_BLUE
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(11)
            p.font.bold = True
            p.font.color.rgb = WHITE
            p.alignment = PP_ALIGN.CENTER
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    for ri, row in enumerate(rows, 1):
        for ci, val in enumerate(row):
            cell = table.cell(ri, ci)
            cell.text = val
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(10)
                p.font.color.rgb = DARK
            if ri % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT
    _footer(slide)


def _closing(prs) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, POSCO_BLUE)
    t = slide.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(11.5), Inches(1.2))
    t.text_frame.text = "Smart Facility, Safe Future"
    t.text_frame.paragraphs[0].font.size = Pt(36)
    t.text_frame.paragraphs[0].font.bold = True
    t.text_frame.paragraphs[0].font.color.rgb = WHITE
    t.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    s = slide.shapes.add_textbox(Inches(1.5), Inches(3.6), Inches(10), Inches(1.5))
    s.text_frame.word_wrap = True
    s.text_frame.text = (
        "경청해 주셔서 감사합니다.\n"
        "질문·개선 제안은 Smart FMS 시스템관리자에게 연락 부탁드립니다."
    )
    for p in s.text_frame.paragraphs:
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(0xCC, 0xDD, 0xEE)
        p.alignment = PP_ALIGN.CENTER
        p.space_after = Pt(8)

    f = slide.shapes.add_textbox(Inches(0.8), Inches(6.5), Inches(11.5), Inches(0.4))
    f.text_frame.text = "POSCO WIDE  ·  Smart FMS"
    f.text_frame.paragraphs[0].font.size = Pt(14)
    f.text_frame.paragraphs[0].font.color.rgb = GOLD
    f.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER


def build() -> Path:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # 1 표지
    _title_cover(
        prs,
        "QR 기반 시설관리\n스마트 AI 통합 플랫폼 구축",
        [
            "POSCO WIDE Smart FMS (Facility Management System)",
            "현장 QR · 예방점검(PM) · 정비 CMMS · 점검일지 · 가로등 · 위험성평가",
            "자체 구축 · 클라우드 운영 · 전체 자료 ZIP 백업",
        ],
    )

    # 2 목차
    _content(
        prs,
        "발표 순서",
        [
            "01  기획 배경 — 시설관리 디지털 전환 로드맵",
            "02  핵심 혁신 ① 설비 QR · 예방점검(PM)",
            "03  핵심 혁신 ② 가로등 · 민원형 현장 접수",
            "04  핵심 혁신 ③ 자산·정비·작업허가 통합",
            "05  구축·적용 현황 (운영 수치)",
            "06  운영 효과 · 기대 성과",
            "07  QR 시연 시나리오 · 비전",
        ],
    )

    # 3 로드맵
    _section(prs, "01  Introduction", "시설 관리 디지털 전환 로드맵", "수기 → 과도기 → Smart FMS 통합 플랫폼")
    _content(
        prs,
        "디지털 전환 로드맵",
        [
            "과거 — 아날로그·수기",
            "  · 종이 점검표·엑셀 분산 관리",
            "  · 이력 추적·통계·누락 파악의 한계",
            "",
            "과도기 — 부분 디지털",
            "  · 건물별 엑셀·공유폴더 중심",
            "  · 보안·권한·실시간 관제 대시보드 부재",
            "",
            "현재 — Smart FMS 통합 플랫폼",
            "  · QR·모바일·웹 원스톱, 역할·메뉴 권한",
            "  · PM·정비·D-1 작업허가·점검일지·가로등·위험성평가",
            "  · 실시간 Dashboard · ZIP 전체 백업 · 확장(API·AI) 준비",
        ],
    )

    # 4 혁신1 설비
    _section(prs, "02  Core Innovation", "설비 QR · 스마트 예방점검", "종이 점검표를 QR·PM으로 대체")
    _two_col(
        prs,
        "설비관리 · 예방점검(PM)",
        "현장 QR (/eq/{코드})",
        [
            "설비 사양·위치·이력을 즉시 조회",
            "PM 결과(정상/주의/고장) 모바일 등록",
            "이상 시 정비의뢰로 바로 연계",
            "점검일지「일지작성」원클릭 연결",
            "로그인 없이도 QR에서 일지 편집·저장",
        ],
        "주요 혁신 포인트",
        [
            "무결성: DB 기반 이력, 임의 수기 조작 차단",
            "누락 방지: PM 주기·금일/지연 KPI Dashboard",
            "엑셀 Import/Export로 기존 설비현황 즉시 이관",
            "건물별 설비현황·설비상세 백업으로 복원 가능",
            "도면·표준서 첨부로 현장 매뉴얼 자산화",
        ],
    )

    # 5 혁신2 가로등/민원
    _section(prs, "03  Customer Experience", "가로등 QR · 스마트 민원 접수", "시민·현장 불편을 QR로 다이렉트 접수")
    _content(
        prs,
        "가로등 · 민원형 접수 혁신",
        [
            "시민(비로그인) — /lamp/{코드}",
            "  · 불점등·글로브파손·전도위험·밝기불량·기타 의뢰",
            "  · 전화 끝 4자리로 본인 확인 후 접수",
            "  · /status 에서 이름+전화로 진행상황 조회",
            "",
            "관리자 — 정비관리 → 가로등",
            "  · 의뢰 목록·상태·메모·Excel 내보내기",
            "  · QR ZIP(전체/구역별) · CSV·엑셀 일괄 등록",
            "  · 일일 리포트 메일·SMS(설정 시)",
            "",
            "기존 타 플랫폼의 ‘공용부 민원 QR’과 같은 경험을,",
            "가로등·설비 QR로 Smart FMS에 내재화했습니다.",
        ],
    )

    # 6 혁신3 자산/정비
    _section(prs, "04  Asset & CMMS", "자산·정비·작업허가 통합", "사후 복구 → 선제 정비·허가 워크플로")
    _content(
        prs,
        "정비 3단계 · 자산 디지털화",
        [
            "① 정비섹션 (/admin/work-orders)",
            "  · 접수 → 진행 → 완료 · 협력사 지정 · D-1 승인",
            "② D-1 / 협력사 (/admin/d1)",
            "  · 예정일·조치·진행 · 작업허가 승인요청",
            "③ 시설섹션 (/admin/facility-section)",
            "  · 잠재위험·안전대책 확인 후 작업허가 최종 승인",
            "",
            "자산·문서",
            "  · 사업장→건물→층→구역→설비 계층 DB",
            "  · 점검일지 엑셀(OnlyOffice/간단편집) · 도면·표준서",
            "  · 위험성평가(로컬/AI) · 자재·협력사 마스터",
            "  · 완료 시 설비 정비이력 자동 등록",
        ],
    )

    # 7 현황 수치
    _section(prs, "05  Operations Status", "구축 · 적용 현황", "현재 Smart FMS에 적재된 운영 규모")
    _metrics(
        prs,
        "Smart FMS 운영 규모 (로컬/운영 DB 기준)",
        [
            ("40", "건물", "사업장 계층으로\n설비·일지·도면 관리"),
            ("3,227", "활성 설비", "건물별 설비현황\nQR·엑셀·상세 백업"),
            ("1,157", "가로등", "시민 QR 의뢰\n· 진행 조회"),
            ("통합", "모듈", "PM·정비·D-1\n점검일지·위험성평가"),
        ],
    )
    _table(
        prs,
        "모듈별 적용 상태",
        ["영역", "상태", "비고"],
        [
            ["통합 플랫폼 (웹·모바일·PWA)", "상시 가동", "역할·메뉴 권한, Dashboard KPI"],
            ["설비 QR · PM · 점검일지", "전면 운영", "/eq/{코드}, /eq/{코드}/log"],
            ["정비 CMMS · D-1 · 작업허가", "전면 운영", "3단계 워크플로"],
            ["가로등 시민 QR", "구축 완료", "/lamp/{코드}, /status"],
            ["위험성평가 · 자재 · 협력사", "운영", "AI 모드(계정별 키) 선택"],
            ["전체 ZIP 백업", "운영", "설비현황·상세·업로드 원본 포함"],
        ],
    )

    # 8 효과
    _section(prs, "06  Key Metrics", "운영 효과 · 기대 성과", "개별 QR 솔루션을 넘어 건물·전사 통합 관제")
    _two_col(
        prs,
        "기대 효과",
        "현장 · 운영",
        [
            "점검·일지·정비 이력의 단일 DB화",
            "QR로 현장→시스템 다이렉트 연결",
            "D-1·작업허가로 안전 프로세스 정착",
            "Dashboard로 지연·긴급 즉시 파악",
            "백업으로 데이터 유실·이전 리스크 감소",
        ],
        "확장 · 거버넌스",
        [
            "역할별 CRUD·메뉴 접근 세분화",
            "시스템관리자 계정·서버 상태 관제",
            "API·AI CCTV 등 확장 인터페이스 준비",
            "위험성평가 AI·보고서 고도화 가능",
            "구독형 외주 의존 없이 자체 플랫폼 보유",
        ],
    )

    # 9 시연
    _section(prs, "07  Live Demo", "QR 스캔 시연 시나리오", "발표장에서 바로 보여줄 URL")
    _table(
        prs,
        "실시간 시연 URL",
        ["시나리오", "URL", "포인트"],
        [
            ["설비 조회·PM", "/eq/{설비코드}", "사양·이력·점검 등록"],
            ["점검일지 작성", "/eq/{코드}/log", "연결 엑셀 즉시 편집"],
            ["가로등 의뢰", "/lamp/{가로등코드}", "시민 비로그인 접수"],
            ["진행 조회", "/status", "이름+전화 4자리"],
            ["관리자 Dashboard", "/admin/dashboard", "KPI·딥링크"],
            ["ZIP 백업", "/admin/server", "전체 자료 내려받기"],
        ],
    )
    _content(
        prs,
        "플랫폼 지향 철학",
        [
            "“측정할 수 없으면 관리할 수 없고,",
            " 관리할 수 없으면 개선할 수 없다.”",
            "  — 피터 드러커 (Peter Drucker)",
            "",
            "Smart FMS는 현장의 모든 점검·정비·민원 행동을",
            "데이터로 측정·자산화하여,",
            "가장 안전한 스마트 시설 환경을 지속적으로 만들겠습니다.",
        ],
    )

    # closing
    _closing(prs)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    return OUT


if __name__ == "__main__":
    path = build()
    print(f"OK {path} {path.stat().st_size // 1024}KB")
