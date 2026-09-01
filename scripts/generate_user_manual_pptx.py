#!/usr/bin/env python3
"""Smart FMS 사용자 매뉴얼 PPT 생성 스크립트."""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parent.parent
OUTPUT = ROOT / "docs" / "Smart_FMS_사용매뉴얼.pptx"

# POSCO WIDE Smart FMS 톤
COLOR_PRIMARY = RGBColor(0x00, 0x3D, 0x7A)
COLOR_ACCENT = RGBColor(0x00, 0x7A, 0xCC)
COLOR_TEXT = RGBColor(0x22, 0x22, 0x22)
COLOR_MUTED = RGBColor(0x66, 0x66, 0x66)
FONT_TITLE = "맑은 고딕"
FONT_BODY = "맑은 고딕"

SLIDES: list[tuple[str, list[str]]] = [
    (
        "Smart FMS\n사용자 매뉴얼",
        [
            "POSCO WIDE 스마트 시설관리 플랫폼",
            "Facility Management System",
            "예방점검 · 정비(CMMS) · 점검일지 · QR 현장운영 · AI 분석",
            "POSCO WIDE Smart FMS  ·  2026.09",
        ],
    ),
    (
        "목차",
        [
            "01  시스템 개요 · 로그인 · 역할",
            "02  메뉴 구성 · 권한",
            "03  Dashboard · 사업장/건물/설비",
            "04  점검(PM) · 점검일지1 · 점검일지2",
            "05  정비 프로세스 (정비섹션 · D-1 · 시설섹션)",
            "06  자재 · 협력사 · 위험성평가 · 가로등",
            "07  AI 분석 · 서버관리 · QR URL 부록",
        ],
    ),
    (
        "01  시스템 개요",
        [
            "Smart FMS는 사업장·건물·설비 기준정보, 예방점검(PM),",
            "점검일지(엑셀·웹), 정비의뢰·D-1·작업허가, 자재·협력사,",
            "위험성평가, AI 분석, 가로등 QR 의뢰를 하나의 웹에서 운영합니다.",
            "PC 브라우저 중심이며, QR 스캔으로 모바일 현장 입력을 지원합니다.",
            "배포: Render (main 브랜치 자동 배포) · DB: PostgreSQL/SQLite",
        ],
    ),
    (
        "로그인 · 가입 · 내 계정",
        [
            "로그인: /admin/login",
            "회원가입: /admin/signup → 관리자 승인 후 사용 가능",
            "내 계정: /admin/account (모든 로그인 사용자 접근)",
            "  · OpenAI API 키·모델 설정 (AI 상세분석·대화·위험성평가)",
            "  · 비밀번호 변경",
            "최초 관리자: 환경변수 ADMIN_ID / ADMIN_PW (배포 시 변경 권장)",
        ],
    ),
    (
        "02  메뉴 구성 (PC 사이드바)",
        [
            "Dashboard · 주요설비 일정 · 공지사항",
            "사업장/건물 · 설비관리 · 점검(PM)",
            "점검일지 · 점검일지2",
            "정비접수/승인(정비섹션) · 정비 List(D-1)/협력사",
            "작업허가/승인(시설섹션) · 가로등 · AI 분석",
            "서버관리 · 위험성평가 · 자재관리 · 협력사 · 계정관리",
            "메뉴는 계정별 menu_access로 개별 허용·차단 가능",
        ],
    ),
    (
        "역할 · 권한",
        [
            "역할: 시스템관리자 · 사업장관리자 · 그룹장 · 파트장 ·",
            "      시설담당자 · 협력사 · 외부업체 · 조회전용",
            "CRUD: can_create / can_edit / can_delete (역할 기본값 + 계정별 조정)",
            "시스템관리자 전용: 서버관리, 계정관리",
            "협력사·외부업체: 설비/PM/점검일지/시설섹션/가로등 메뉴 기본 차단",
            "조회전용: 읽기만 가능 (추가·수정·삭제 불가)",
        ],
    ),
    (
        "03  Dashboard (운영 현황)",
        [
            "경로: /admin/dashboard",
            "KPI: 사업장·건물·설비·PM 지연·정비의뢰·점검 현황 등",
            "위젯: 정비현황 · 사업장현황 · 에너지 · 오늘 일정 · 공지",
            "위젯 순서·표시 여부·레이아웃(gallery/bento) 사용자 설정 가능",
            "30초 간격 자동 갱신 (캐시 TTL)",
        ],
    ),
    (
        "04  사업장 · 건물 · 설비",
        [
            "사업장/건물: /admin/sites — 사업장·건물·층·구역 CRUD",
            "  · 도면·표준서 파일 업로드",
            "설비관리: /admin/equipment — 설비 등록·수정·삭제",
            "  · 사업장·건물별 트리 탐색 · 상세 스펙·이력",
            "  · 엑셀 일괄 Import / Export (전체·건물별)",
            "  · QR 코드 생성 (개별 PNG · 건물별 ZIP)",
            "설비에서 정비의뢰·PM·점검일지 바로 연계",
        ],
    ),
    (
        "설비 QR · 현장 화면",
        [
            "설비 QR URL: {BASE}/eq/{설비코드}",
            "  · 모바일에서 설비 상세·정비의뢰·PM 점검·일지 작성",
            "점검일지(로그인 불필요): {BASE}/eq/{코드}/log",
            "QR PNG: /admin/equipment/{id}/qr.png",
            "건물 일괄 QR ZIP: /admin/equipment/building/{building_id}/qr-zip",
            "PUBLIC_BASE_URL 환경변수로 QR 도메인 고정 (배포 URL과 일치 권장)",
        ],
    ),
    (
        "05  점검(PM) — 예방점검",
        [
            "경로: /admin/pm",
            "점검 일정 등록: 매일/매주/매월/분기/반기/연간/사용자정의",
            "점검 실행 → 결과: 정상 / 주의 / 고장",
            "지연·기한 필터 · 엑셀 Export",
            "QR 현장 점검: /eq/{코드} → PM 점검 (고장 시 정비의뢰 연계)",
        ],
    ),
    (
        "06  점검일지1 — 엑셀 기반",
        [
            "경로: /admin/inspection-logs",
            "건물별 점검일지 등록 → 엑셀 파일 업로드",
            "OnlyOffice 또는 레거시 편집기로 웹 편집",
            "QR 연결 일지 작성: {BASE}/eq/{코드}/log",
            "건물 추가·삭제 시 사이드바 메뉴 즉시 반영",
        ],
    ),
    (
        "07  점검일지2 — 웹 운영일보 개요",
        [
            "경로: /admin/inspection-logs2",
            "건물명으로 유형 자동 판별 → 전용 웹 화면",
            "1일 입력 · 일 마감(close-day) · 월보/년보 엑셀 Export",
            "QR로 현장 1일 입력 (로그인 불필요)",
            "마감 시 엑셀 아카이브 생성·다운로드",
            "유형: 주택변전소 · 중앙관제실 · 중앙관제실(설비) · 제철소본부",
        ],
    ),
    (
        "점검일지2 — 주택변전소",
        [
            "판별: 건물명에 '주택변전소' 포함",
            "관리: /admin/inspection-logs2/{id}/housing",
            "QR 1일 입력: {BASE}/hs/{건물코드}/daily",
            "전력·운전 데이터 자동 집계 · 월보/년보 kWh",
            "QR PNG: .../housing/qr.png",
        ],
    ),
    (
        "점검일지2 — 중앙관제실",
        [
            "판별: 건물명 '중앙관제실' (설비 유형 제외)",
            "관리: /admin/inspection-logs2/{id}/central-control-room",
            "QR 1일 입력: {BASE}/ccr/{건물코드}/daily",
            "일·월·년 엑셀 Export · 아카이브 다운로드",
            "QR PNG: .../central-control-room/qr.png",
        ],
    ),
    (
        "점검일지2 — 중앙관제실(설비)",
        [
            "판별: 건물명 '중앙관제실(설비)'",
            "관리: /admin/inspection-logs2/{id}/ccr-facility",
            "QR 1일 입력: {BASE}/ccrf/{건물코드}/daily",
            "설비별 운전·점검 항목 웹 입력 · daily Export",
            "QR PNG: .../ccr-facility/qr.png",
        ],
    ),
    (
        "점검일지2 — 제철소본부",
        [
            "판별: 건물명 '제철소본부'",
            "관리: /admin/inspection-logs2/{id}/steelworks-hq",
            "QR 1일 입력: {BASE}/swhq/{건물코드}/daily",
            "NO.1/NO.2 수전 배율 12000 · 일사용량 kWh 자동 계산",
            "전일 월누계 기반 일일·월누계 재계산",
            "QR PNG: .../steelworks-hq/qr.png",
        ],
    ),
    (
        "08  정비 프로세스 개요",
        [
            "설비 QR·관리 화면·수동 접수 → 정비의뢰(received)",
            "  → 배정(assigned) → D-1 승인 → D-1 계획 보드",
            "  → 작업허가 요청 → 시설섹션 승인(permit)",
            "  → in_progress → completed → verified → closed",
            "정비섹션·D-1·시설섹션 메뉴에 신규 건수 배지 표시",
        ],
    ),
    (
        "정비접수/승인 (정비섹션)",
        [
            "경로: /admin/work-orders",
            "정비의뢰 목록·상태 필터·상세 조회",
            "D-1 승인 · 일괄 D-1 승인",
            "작업허가 요청(request-approval) → 시설섹션으로 전달",
            "엑셀 Export",
        ],
    ),
    (
        "정비 List(D-1) / 협력사",
        [
            "경로: /admin/d1",
            "D-1 계획 보드 · 협력사별 작업 필터",
            "단계 진행: draft → review → approved → JSA → TBM",
            "  → permit → 진행 → 완료",
            "협력사 계정과 partner_id 연계",
        ],
    ),
    (
        "작업허가/승인 (시설섹션)",
        [
            "경로: /admin/facility-section",
            "보드: 전일(day_before) · 당일(today) · 예정(scheduled) · 전체",
            "작업허가 승인 · 일괄 승인",
            "협력사별 위험등급·기본값 설정",
        ],
    ),
    (
        "09  자재관리",
        [
            "경로: /admin/materials (팝업: ?popup=1)",
            "품목 등록 · 입고(stock-in) · 출고(stock-out)",
            "그룹 관리 · 재고 로그",
            "엑셀 Import / Export · 초기화(reset)",
            "설비 화면에서 팝업으로 연동 가능",
        ],
    ),
    (
        "협력사 관리",
        [
            "경로: /admin/partners",
            "협력사 코드·담당자·계약종료일 관리",
            "정비의뢰 D-1·시설섹션에서 partner_id 필터",
            "협력사·외부업체 계정: partner_id 또는 company_name 연결",
        ],
    ),
    (
        "위험성평가",
        [
            "경로: /admin/risk-assessment",
            "작업명 + 4M+1E(인·기계·재료·방법·환경) 입력",
            "대분류·프리셋 선택 · AI 보조 평가 (OpenAI API 키 필요)",
            "과거 평가 문서 업로드·학습으로 사내 기준 반영",
            "HTML / Excel 보고서 Export",
        ],
    ),
    (
        "10  AI 분석",
        [
            "경로: /admin/ai-analysis",
            "집계 Q&A: DB 스냅샷 기반 즉시 답변 (API 키 불필요)",
            "GPT 상세 모드: FMS 전체 컨텍스트 + GPT 분석",
            "데이터 범위: 사업장·설비·정비·PM·D-1·점검일지1/2·",
            "  자재·공지·일정·가로등·주택변전소 월보 전력 등",
            "API 설정: 내 계정 또는 AI 분석 화면에서 OpenAI 키·모델",
        ],
    ),
    (
        "AI 대화형 채팅",
        [
            "경로: /admin/ai-analysis (채팅 패널)",
            "API: POST /admin/ai-analysis/chat",
            "세션별 최대 20턴 대화 이력 유지",
            "질문 의도 자동 분류 → 관련 DB 섹션 강화 조회",
            "GPT 응답 + 집계 근거(evidence) 병행 표시",
            "OpenAI API 키 미설정 시 버튼 활성 + 안내 메시지",
        ],
    ),
    (
        "가로등",
        [
            "공개 의뢰: {BASE}/lamp/{가로등코드}",
            "  · 정비 유형 선택 → 접수 · /status 로 조회",
            "관리: /admin/streetlamp/requests — 목록·상태·삭제·Export",
            "QR 일괄: /admin/streetlamp/qr-zip",
            "설정: SMS/메일 알림 · CSV import · 일일 리포트",
        ],
    ),
    (
        "주요설비 일정 · 공지사항",
        [
            "일정: /admin/schedules — 캘린더 등록·삭제",
            "  · 카테고리: 긴급 / 점검 / 작업 / 검수",
            "공지: /admin/notices — 등록·삭제",
            "  · 카테고리: 긴급 / 안전 / 일반",
            "Dashboard 위젯과 연동",
        ],
    ),
    (
        "11  QR · URL 정리 (부록)",
        [
            "설비: {BASE}/eq/{코드}  ·  일지: .../eq/{코드}/log",
            "주택변전소: {BASE}/hs/{코드}/daily",
            "중앙관제실: {BASE}/ccr/{코드}/daily",
            "중앙관제실(설비): {BASE}/ccrf/{코드}/daily",
            "제철소본부: {BASE}/swhq/{코드}/daily",
            "가로등: {BASE}/lamp/{코드}",
            "{BASE} = PUBLIC_BASE_URL 또는 사이트 도메인",
        ],
    ),
    (
        "12  서버관리 (시스템관리자)",
        [
            "경로: /admin/server",
            "서버 상태: Render 배포·DB·디스크·메모리·CPU (30초 갱신)",
            "전체 백업 ZIP: /admin/server/backup.zip",
            "  · 점검일지·도면·표준서 + DB 업무 데이터 엑셀",
            "사용 매뉴얼 PPT: /admin/server/manual.pptx",
            "발표자료 PPT: /admin/server/presentation.pptx",
        ],
    ),
    (
        "계정관리 (시스템관리자)",
        [
            "경로: /admin/users",
            "가입 승인 · 역할 지정 · CRUD 권한 조정",
            "메뉴 접근(menu_access) 개별 설정",
            "협력사·외부업체 partner_id / company_name 연결",
        ],
    ),
    (
        "부록 — 점검일지2 유형·URL",
        [
            "주택변전소     → /hs/{code}/daily     → .../housing",
            "중앙관제실     → /ccr/{code}/daily    → .../central-control-room",
            "중앙관제실(설비) → /ccrf/{code}/daily → .../ccr-facility",
            "제철소본부     → /swhq/{code}/daily   → .../steelworks-hq",
            "리다이렉트 우선순위: 주택 → 제철소 → CCR(설비) → CCR",
        ],
    ),
    (
        "감사합니다",
        [
            "Smart FMS 사용자 매뉴얼",
            "문의: 시스템관리자 · POSCO WIDE Smart FMS",
            "서버관리 → 「사용 매뉴얼 PPT 받기」에서 최신본 다운로드",
        ],
    ),
]


def _set_run_font(run, *, size: int, bold: bool = False, color: RGBColor = COLOR_TEXT):
    run.font.name = FONT_BODY
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def _add_title(slide, text: str, *, cover: bool = False):
    left, top, width, height = Inches(0.6), Inches(0.45), Inches(12.1), Inches(1.2)
    if cover:
        top = Inches(2.0)
        height = Inches(2.0)
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT if not cover else PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    _set_run_font(run, size=32 if cover else 24, bold=True, color=COLOR_PRIMARY)


def _add_bullets(slide, lines: list[str], *, cover: bool = False):
    top = Inches(3.6) if cover else Inches(1.55)
    height = Inches(3.8) if cover else Inches(5.5)
    box = slide.shapes.add_textbox(Inches(0.75), top, Inches(12.0), height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER if cover else PP_ALIGN.LEFT
        p.space_after = Pt(6)
        p.level = 0
        text = line.strip()
        if not cover and not text.startswith("  ") and text and not text.startswith("·"):
            text = "• " + text
        run = p.add_run()
        run.text = text
        size = 16 if cover else (14 if text.startswith("  ") or text.startswith("·") else 15)
        _set_run_font(
            run,
            size=size,
            bold=cover and i == 0,
            color=COLOR_MUTED if cover and i > 0 else COLOR_TEXT,
        )


def _add_footer(slide, page: int, total: int):
    box = slide.shapes.add_textbox(Inches(10.5), Inches(7.0), Inches(2.3), Inches(0.35))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = f"{page} / {total}"
    _set_run_font(run, size=10, color=COLOR_MUTED)


def build_presentation() -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    total = len(SLIDES)

    for idx, (title, bullets) in enumerate(SLIDES, start=1):
        slide = prs.slides.add_slide(blank)
        cover = idx == 1
        _add_title(slide, title, cover=cover)
        _add_bullets(slide, bullets, cover=cover)
        if not cover:
            _add_footer(slide, idx, total)

        # 상단 accent bar
        bar = slide.shapes.add_shape(
            1, Inches(0), Inches(0), prs.slide_width, Inches(0.12)
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = COLOR_ACCENT
        bar.line.fill.background()

    return prs


def main() -> None:
    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs = build_presentation()
    prs.save(str(OUTPUT))
    print(f"Wrote {OUTPUT} ({len(SLIDES)} slides)")


if __name__ == "__main__":
    main()
