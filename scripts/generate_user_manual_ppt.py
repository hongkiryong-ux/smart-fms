# -*- coding: utf-8 -*-
"""Smart FMS 사용 매뉴얼 PPT 생성.

원칙:
- 섹션 제목 전용 슬라이드를 두지 않고, 제목 + 상세 설명을 한 페이지에 배치
- 기능별 절차·주의사항·URL을 자세히 기술하되 페이지 수는 최소화
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "docs" / "Smart_FMS_사용매뉴얼.pptx"

POSCO_BLUE = RGBColor(0x00, 0x38, 0x76)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x33, 0x33, 0x33)
GRAY = RGBColor(0x66, 0x66, 0x66)
LIGHT_BG = RGBColor(0xF0, 0xF4, 0xF8)
ACCENT = RGBColor(0x00, 0x6E, 0xB8)


def _set_slide_bg(slide, color: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_bar(slide, prs) -> None:
    shape = slide.shapes.add_shape(
        1,
        Inches(0),
        Inches(0),
        prs.slide_width,
        Inches(0.42),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = POSCO_BLUE
    shape.line.fill.background()


def _footer(slide, text: str = "POSCO WIDE  ·  Smart FMS 사용 매뉴얼") -> None:
    box = slide.shapes.add_textbox(Inches(0.55), Inches(7.1), Inches(12.2), Inches(0.3))
    p = box.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(10)
    p.font.color.rgb = GRAY


def _fill_lines(tf, lines: list[str], *, base_size: int = 14) -> None:
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.space_after = Pt(4)
        p.level = 0
        if not line.strip():
            p.font.size = Pt(6)
            continue
        if line.startswith("  ·") or line.startswith("    ·"):
            p.level = 1
            p.font.size = Pt(base_size - 1)
            p.font.color.rgb = GRAY
        elif line.startswith("■"):
            p.font.size = Pt(base_size + 1)
            p.font.bold = True
            p.font.color.rgb = ACCENT
            p.space_before = Pt(6)
        else:
            p.font.size = Pt(base_size)
            p.font.color.rgb = DARK


def _title_slide(prs, title: str, subtitle: str = "") -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, POSCO_BLUE)
    box = slide.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(11.5), Inches(1.4))
    tf = box.text_frame
    tf.text = title
    p = tf.paragraphs[0]
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = WHITE
    if subtitle:
        sub = slide.shapes.add_textbox(Inches(0.8), Inches(3.6), Inches(11.5), Inches(1.8))
        stf = sub.text_frame
        stf.word_wrap = True
        for i, line in enumerate(subtitle.split("\n")):
            sp = stf.paragraphs[0] if i == 0 else stf.add_paragraph()
            sp.text = line
            sp.font.size = Pt(18)
            sp.font.color.rgb = RGBColor(0xCC, 0xDD, 0xEE)
            sp.space_after = Pt(6)
    foot = slide.shapes.add_textbox(Inches(0.8), Inches(6.5), Inches(11), Inches(0.5))
    foot.text_frame.text = "POSCO WIDE Smart FMS  ·  2026"
    foot.text_frame.paragraphs[0].font.size = Pt(14)
    foot.text_frame.paragraphs[0].font.color.rgb = RGBColor(0xAA, 0xBB, 0xCC)


def _content_slide(
    prs,
    title: str,
    bullets: list[str],
    note: str = "",
    *,
    base_size: int = 14,
) -> None:
    """제목 + 본문 설명을 한 페이지에 배치."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, WHITE)
    _add_bar(slide, prs)
    tbox = slide.shapes.add_textbox(Inches(0.55), Inches(0.55), Inches(12.2), Inches(0.55))
    tbox.text_frame.text = title
    tbox.text_frame.paragraphs[0].font.size = Pt(22)
    tbox.text_frame.paragraphs[0].font.bold = True
    tbox.text_frame.paragraphs[0].font.color.rgb = POSCO_BLUE

    body_h = 5.7 if not note else 5.35
    body = slide.shapes.add_textbox(Inches(0.55), Inches(1.15), Inches(12.2), Inches(body_h))
    _fill_lines(body.text_frame, bullets, base_size=base_size)

    if note:
        nbox = slide.shapes.add_textbox(Inches(0.55), Inches(6.55), Inches(12.2), Inches(0.45))
        nbox.text_frame.text = f"※ {note}"
        nbox.text_frame.paragraphs[0].font.size = Pt(11)
        nbox.text_frame.paragraphs[0].font.color.rgb = GRAY
    _footer(slide)


def _two_col_slide(
    prs,
    title: str,
    left_title: str,
    left: list[str],
    right_title: str,
    right: list[str],
    note: str = "",
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, WHITE)
    _add_bar(slide, prs)
    tbox = slide.shapes.add_textbox(Inches(0.55), Inches(0.55), Inches(12.2), Inches(0.5))
    tbox.text_frame.text = title
    tbox.text_frame.paragraphs[0].font.size = Pt(22)
    tbox.text_frame.paragraphs[0].font.bold = True
    tbox.text_frame.paragraphs[0].font.color.rgb = POSCO_BLUE

    for col_x, col_title, items in [
        (0.55, left_title, left),
        (6.85, right_title, right),
    ]:
        h = slide.shapes.add_textbox(Inches(col_x), Inches(1.15), Inches(5.7), Inches(0.38))
        h.text_frame.text = col_title
        h.text_frame.paragraphs[0].font.size = Pt(15)
        h.text_frame.paragraphs[0].font.bold = True
        h.text_frame.paragraphs[0].font.color.rgb = ACCENT
        body = slide.shapes.add_textbox(Inches(col_x), Inches(1.55), Inches(5.7), Inches(5.0))
        _fill_lines(body.text_frame, items, base_size=12)

    if note:
        nbox = slide.shapes.add_textbox(Inches(0.55), Inches(6.6), Inches(12.2), Inches(0.4))
        nbox.text_frame.text = f"※ {note}"
        nbox.text_frame.paragraphs[0].font.size = Pt(11)
        nbox.text_frame.paragraphs[0].font.color.rgb = GRAY
    _footer(slide)


def _table_slide(
    prs,
    title: str,
    headers: list[str],
    rows: list[list[str]],
    note: str = "",
    *,
    intro: list[str] | None = None,
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, WHITE)
    _add_bar(slide, prs)
    tbox = slide.shapes.add_textbox(Inches(0.55), Inches(0.55), Inches(12.2), Inches(0.45))
    tbox.text_frame.text = title
    tbox.text_frame.paragraphs[0].font.size = Pt(22)
    tbox.text_frame.paragraphs[0].font.bold = True
    tbox.text_frame.paragraphs[0].font.color.rgb = POSCO_BLUE

    top = 1.1
    if intro:
        ibox = slide.shapes.add_textbox(Inches(0.55), Inches(1.05), Inches(12.2), Inches(0.7))
        _fill_lines(ibox.text_frame, intro, base_size=12)
        top = 1.75

    cols = len(headers)
    row_h = min(0.38, 4.8 / max(len(rows) + 1, 1))
    table = slide.shapes.add_table(
        len(rows) + 1,
        cols,
        Inches(0.5),
        Inches(top),
        Inches(12.3),
        Inches(row_h * (len(rows) + 1)),
    ).table

    for ci, h in enumerate(headers):
        cell = table.cell(0, ci)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = POSCO_BLUE
        for p in cell.text_frame.paragraphs:
            p.font.bold = True
            p.font.size = Pt(11)
            p.font.color.rgb = WHITE
            p.alignment = PP_ALIGN.CENTER
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    for ri, row in enumerate(rows, start=1):
        for ci, val in enumerate(row):
            cell = table.cell(ri, ci)
            cell.text = val
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(10)
                p.font.color.rgb = DARK
            if ri % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT_BG

    if note:
        nbox = slide.shapes.add_textbox(Inches(0.55), Inches(6.6), Inches(12.2), Inches(0.4))
        nbox.text_frame.text = f"※ {note}"
        nbox.text_frame.paragraphs[0].font.size = Pt(11)
        nbox.text_frame.paragraphs[0].font.color.rgb = GRAY
    _footer(slide)


def build() -> Path:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # ── 표지 ──
    _title_slide(
        prs,
        "Smart FMS\n사용 매뉴얼",
        "POSCO WIDE 통합 시설관리 플랫폼\nFacility Management System\n예방점검 · 정비(CMMS) · 점검일지 · QR 현장운영",
    )

    # ── 목차 (제목+설명 한 장) ──
    _content_slide(
        prs,
        "목차 · 이 매뉴얼의 구성",
        [
            "■ 01 시스템 개요·접속 — Smart FMS 범위, 로그인/가입, 화면 구성",
            "■ 02 메뉴·권한 — PC/모바일 메뉴, 역할별 CRUD·메뉴 접근",
            "■ 03 Dashboard — KPI 위젯, 레이아웃·설정, 카드 클릭 이동",
            "■ 04 사업장·건물·설비 — 계층 구조, 도면/표준서, QR",
            "■ 05 엑셀 Import/Export — 건물별·일괄(폴더) 대체 Import, Export",
            "■ 06 점검(PM)·점검일지 — 주기/결과, OnlyOffice·편집 잠금, QR 일지",
            "■ 07 정비관리 3단계 — 접수/승인 → D-1 → 작업허가/승인",
            "■ 08 가로등·위험성평가·자재·협력사 — 부가 모듈 운영",
            "■ 09 QR 현장·계정·백업 — 공개 URL, 서버 ZIP 백업, 빠른 참조",
            "",
            "읽는 방법: 각 장은 「제목 + 상세 절차」가 한 페이지에 정리되어 있습니다.",
            "화면 경로는 /admin/... 형식이며, QR·가로등 일부는 로그인 없이 사용합니다.",
        ],
        base_size=13,
    )

    # ── 01 개요 ──
    _content_slide(
        prs,
        "01 · 시스템 개요 및 접속",
        [
            "■ Smart FMS란?",
            "사업장→건물→층→구역→설비를 계층 관리하는 통합 시설관리(FMS)입니다.",
            "예방점검(PM), 정비(CMMS), 점검일지(엑셀), D-1 작업허가, 가로등 QR, 자재·협력사를 한 플랫폼에서 운영합니다.",
            "",
            "■ 접속·계정",
            "  · 관리자 로그인: /admin/login  →  성공 시 /admin/dashboard",
            "  · 가입 신청: /admin/signup (아이디·이름·역할·회사·비밀번호) → 시스템관리자 승인 후 사용",
            "  · 내 계정: /admin/account (본인 정보·비밀번호 변경)",
            "  · 최초 배포 시 admin 계정 자동 생성 (비밀번호: 서버 환경변수 ADMIN_PW)",
            "",
            "■ 화면 구성 요약",
            "  · PC: 좌측 사이드바 전체 메뉴 · 설비/점검일지는 건물 목록 펼침",
            "  · 모바일(약 900px 이하): 햄버거 메뉴 · QR 설비·일지작성이 주 사용 환경 · PWA 홈 추가 가능",
            "  · 자재관리는 메뉴 클릭 시 별도 팝업 창으로 열립니다.",
        ],
        note="문의·권한 문제는 시스템관리자에게 요청하세요.",
        base_size=13,
    )

    # ── 02 메뉴 ──
    _table_slide(
        prs,
        "02 · 메뉴 구조 (PC 사이드바)",
        ["메뉴", "경로", "주요 기능"],
        [
            ["Dashboard", "/admin/dashboard", "KPI·사업장 현황·위젯 설정"],
            ["사업장/건물", "/admin/sites", "사업장·건물·층·구역·도면·표준서"],
            ["설비관리", "/admin/equipment", "분류별 설비·엑셀·QR·정비의뢰"],
            ["점검(PM)", "/admin/pm", "예방점검 주기·결과·Excel"],
            ["점검일지", "/admin/inspection-logs", "건물별 엑셀 일지·편집·연결"],
            ["정비접수/승인", "/admin/work-orders", "정비섹션 CMMS 접수·진행·완료"],
            ["정비 List(D-1)", "/admin/d1", "협력사 작업·허가 요청"],
            ["작업허가/승인", "/admin/facility-section", "시설섹션 최종 작업허가"],
            ["가로등", "/admin/streetlamp/requests", "QR 가로등 정비·리포트"],
            ["위험성평가", "/admin/risk-assessment", "JSA·5M1E·HTML/Excel"],
            ["자재·협력사", "/admin/materials · /partners", "재고·입출고 · 협력사 마스터"],
            ["계정·서버", "/admin/users · /server", "승인·권한 · 상태·ZIP 백업"],
        ],
        intro=[
            "모바일에서는 정비관리 하위 메뉴 중심이며, Dashboard·설비 등은 URL 직접 접근도 가능합니다.",
        ],
        note="메뉴 표시는 역할·개별 메뉴 권한에 따라 달라질 수 있습니다.",
    )

    _table_slide(
        prs,
        "02 · 사용자 역할 및 권한 설정",
        ["역할", "기본 CRUD", "접근 범위 · 비고"],
        [
            ["시스템관리자", "전체", "모든 메뉴 + 계정관리·서버관리"],
            ["사업장관리자", "추가·수정·삭제", "계정·서버 제외 · 운영 전반"],
            ["그룹장 / 파트장", "추가·수정·삭제", "계정·서버 제외"],
            ["시설담당자", "추가·수정·삭제", "시설·정비·점검 중심"],
            ["협력사 / 외부업체", "수정 위주", "정비·D-1·작업 관련 중심"],
            ["조회전용", "조회만", "추가·수정·삭제 불가"],
        ],
        intro=[
            "가입(/admin/signup) → 계정관리에서 승인/거절 → 계정별 CRUD·메뉴 접근을 개별 지정할 수 있습니다.",
            "본인 계정의 역할 변경·삭제는 불가합니다. 계정·서버 메뉴는 시스템관리자 전용입니다.",
        ],
    )

    # ── 03 Dashboard ──
    _content_slide(
        prs,
        "03 · Dashboard (운영 현황)",
        [
            "■ 무엇을 보나?",
            "사업장·건물·설비 수, PM(금일 계획·완료·지연), 정비(의뢰·미해결·완료), D-1 등 KPI를 한눈에 확인합니다.",
            "사업장 현황 카드의 지표를 클릭하면 해당 정비 목록(필터 적용)으로 이동합니다.",
            "",
            "■ 사용 방법",
            "  · 경로: /admin/dashboard  ·  약 30초 주기 자동 갱신",
            "  · 레이아웃: /admin/dashboard/layouts — 운영(ops) / 한눈에(bento) / 시설(gallery)",
            "  · 설정: /admin/dashboard/settings — 위젯 표시·순서(드래그) 조정 후 저장",
            "",
            "■ 팁",
            "KPI 카드·사업장 행을 눌러 하위 화면으로 바로 들어가 조치(정비 접수·점검)를 이어가면 효율적입니다.",
            "권한에 따라 일부 위젯·메뉴가 숨겨질 수 있습니다.",
        ],
        base_size=13,
    )

    # ── 04 사업장·설비 ──
    _two_col_slide(
        prs,
        "04 · 사업장 · 건물 · 설비 관리",
        "사업장 / 건물 (/admin/sites)",
        [
            "■ 계층",
            "사업장 → 건물 → 층 → 구역 → 설비",
            "",
            "■ 주요 작업",
            "  · 사업장·건물 등록·수정·삭제",
            "  · 같은 사업장 내 건물명·코드 중복 불가",
            "  · 건물 상세: 층·구역 추가",
            "  · 도면·표준서 첨부 (파일당 최대 약 20MB)",
            "",
            "■ 주의",
            "건물 삭제·비활성 시 하위 설비·일지 연계를 확인하세요.",
        ],
        "설비관리 (/admin/equipment)",
        [
            "■ 건물 선택",
            "사이드바 또는 건물 카드에서 건물 선택",
            "",
            "■ 기능",
            "  · 분류(엑셀 시트명)별 설비 목록",
            "  · 상세: 사양·정비이력·점검·의뢰",
            "  · 정비의뢰 접수, 변경 로그",
            "  · QR PNG / ZIP 다운로드 → 현장 부착",
            "  · 단일 설비 Excel 내보내기·가져오기",
            "",
            "■ 상세 팝업",
            "목록에서 설비 행을 열면 이력·의뢰를 바로 확인·처리할 수 있습니다.",
        ],
        note="설비 템플릿(/admin/templates)은 URL 직접 접근으로 종류·점검항목을 관리합니다.",
    )

    # ── 05 엑셀 ──
    _content_slide(
        prs,
        "05 · 설비 엑셀 Import / Export (대체 방식)",
        [
            "■ 공통 경고 (Import 시 반드시 확인)",
            "승인하면 해당 범위의 기존 사양·정비이력·점검이력·점검주기·정비의뢰 등이 사라지고",
            "엑셀에 있는 내용으로 전부 대체됩니다. 진행 전 ZIP 백업 또는 엑셀 Export를 권장합니다.",
            "",
            "■ 건물별 Import (/admin/equipment/import)",
            "  · 건물 선택 → xls/xlsx 업로드 → 경고 확인 → 시트별 설비 + 정비/점검이력 반영",
            "  · 총괄·표지 등 메타 시트는 제외, 시트명 = 설비 분류",
            "",
            "■ 건물 일괄 Import (폴더 선택)",
            "  · 폴더를 선택하면 안의 모든 엑셀(.xls/.xlsx)을 파일명(건물명)에 매칭해 일괄 대체",
            "  · 확인창에 엑셀 개수·대체 안내가 표시됩니다",
            "",
            "■ Export",
            "  · 건물「엑셀 출력」→ {건물명}_설비현황.xlsx (시트 + 정비이력 + 점검이력)",
            "  · 설비「Excel 내보내기」→ 사양·정비·점검·주기·정비의뢰 전체",
        ],
        note="단일 설비 가져오기는 xlsx 권장(내보내기 형식과 동일). 일괄은 폴더 내 엑셀만 처리합니다.",
        base_size=12,
    )

    # ── 06 PM · 점검일지 ──
    _two_col_slide(
        prs,
        "06 · 점검(PM) · 점검일지",
        "예방점검 PM (/admin/pm)",
        [
            "■ 탭",
            "점검목록 / 주기설정",
            "",
            "■ 주기",
            "일·주·월·분기·반기·연·사용자지정",
            "",
            "■ 결과",
            "정상 / 주의 / 고장",
            "이상 시 정비의뢰 연계 가능",
            "",
            "■ 기타",
            "Excel 내보내기",
            "QR(/eq/{코드})에서도 PM 등록",
            "일정·담당자·다음예정일 관리",
        ],
        "점검일지 (/admin/inspection-logs)",
        [
            "■ 준비",
            "1) 건물 등록  2) 엑셀 업로드",
            "3)「이 파일로 연결」(건물당 1개)",
            "",
            "■ 편집",
            "OnlyOffice 또는 브라우저 간단 편집기",
            "경로: .../files/{id}/edit",
            "마지막 편집 셀 위치 저장",
            "",
            "■ 동시 편집 잠금",
            "다른 사용자가 편집 중이면",
            "편집자 이름이 안내되고 대기합니다.",
            "편집 종료·이탈 시 잠금 해제",
            "",
            "■ QR 일지",
            "/eq/{코드}/log → 연결 엑셀 바로 작성",
        ],
        note="OnlyOffice는 docker-compose.onlyoffice.yml 로 구성할 수 있습니다.",
    )

    # ── 07 정비 ──
    _content_slide(
        prs,
        "07 · 정비관리 3단계 워크플로",
        [
            "■ ① 정비접수/승인 (정비섹션)  —  /admin/work-orders",
            "정비의뢰 접수 → 정비진행 → 정비완료. 협력사 지정 후 D-1 일괄/개별 승인.",
            "메뉴 배지로 신규·미처리 건을 확인할 수 있습니다.",
            "",
            "■ ② 정비 List(D-1)/협력사  —  /admin/d1",
            "D-1 승인된 건만 표시. 예정일·조치·진행상태 편집 후 작업허가 승인요청(개별·일괄).",
            "보드 필터: 전체 / 접수 / 오늘 / 내일 / 예정 / 완료 등",
            "",
            "■ ③ 작업허가/승인 (시설섹션)  —  /admin/facility-section",
            "잠재위험·안전대책 확인(협력사 기본값·엑셀 매칭) 후 작업허가 최종 승인.",
            "필터: 하루전 / 당일 / 예정 / 전체",
            "",
            "■ 완료 후",
            "정비완료 시 해당 설비의 정비이력에 자동 등록됩니다. (중복 방지 로직 적용)",
        ],
        note="D-1 승인 상태(d1_approved)가 되어야 협력사 D-1 목록에 노출됩니다.",
        base_size=13,
    )

    # ── 08 부가 모듈 ──
    _two_col_slide(
        prs,
        "08 · 가로등 · 위험성평가 · 자재 · 협력사",
        "가로등 · 위험성평가",
        [
            "■ 가로등 (/admin/streetlamp/requests)",
            "정비의뢰 목록·상태·메모·Excel",
            "QR ZIP(전체/구역별), CSV·엑셀 일괄",
            "일일 리포트 메일·SMS 설정",
            "",
            "■ 시민·현장 (로그인 불필요)",
            "  · /lamp/{코드} 가로등 정비 의뢰",
            "  · /status 이름+전화 4자리로 진행 조회",
            "",
            "■ 위험성평가 (/admin/risk-assessment)",
            "로컬/AI 모드, JSA·5M1E",
            "HTML·Excel 출력",
        ],
        "자재 · 협력사",
        [
            "■ 자재관리 (/admin/materials)",
            "메뉴 클릭 시 팝업 창으로 열림",
            "그룹별 재고, 입출고, Excel",
            "",
            "■ 협력사 (/admin/partners)",
            "업체·담당·연락·계약만료 관리",
            "D-1·작업허가·시설섹션의",
            "잠재위험·안전대책 기본값 연동",
            "시트명=업체명 엑셀 카탈로그 지원",
            "",
            "■ 운영 팁",
            "협력사 기본값을 먼저 맞춰 두면",
            "승인 모달 입력이 빨라집니다.",
        ],
    )

    # ── 09 QR · 관리 ──
    _table_slide(
        prs,
        "09 · QR 현장 URL · 계정 · 서버 백업",
        ["구분", "URL / 경로", "설명"],
        [
            ["설비 QR", "/eq/{설비코드}", "로그인 없이 조회·PM·정비이력"],
            ["일지작성", "/eq/{코드}/log", "연결 엑셀 바로 편집·저장"],
            ["가로등 QR", "/lamp/{코드}", "시민 정비 의뢰"],
            ["진행조회", "/status", "가로등 진행(이름+전화 4자리)"],
            ["계정관리", "/admin/users", "가입 승인·역할·CRUD·메뉴 (관리자)"],
            ["서버관리", "/admin/server", "CPU·메모리·디스크·DB · 30초 갱신"],
            ["ZIP 백업", "/admin/server/backup.zip", "files/ + excel/설비현황·상세 + 업무데이터"],
            ["Health", "/health", "서비스 생존 확인"],
        ],
        intro=[
            "QR 일지: 점검일지에서 건물 등록→엑셀 업로드→「이 파일로 연결」→ 현장 QR「일지작성」.",
            "ZIP 백업은 무압축 스트리밍이며, Import 전·대규모 변경 전 받아 두는 것을 권장합니다.",
        ],
        note="files/는 점검일지·도면·표준서 원본, excel/은 UI Export와 동일 형식입니다.",
    )

    _content_slide(
        prs,
        "09 · URL 빠른 참조 · 문의",
        [
            "■ 기본",
            "로그인 /admin/login  ·  가입 /admin/signup  ·  Dashboard /admin/dashboard  ·  내 계정 /admin/account",
            "",
            "■ 시설·점검",
            "사업장 /admin/sites  ·  설비 /admin/equipment  ·  Import /admin/equipment/import",
            "PM /admin/pm  ·  점검일지 /admin/inspection-logs",
            "",
            "■ 정비",
            "접수/승인 /admin/work-orders  ·  D-1 /admin/d1  ·  시설섹션 /admin/facility-section",
            "",
            "■ 기타",
            "가로등 /admin/streetlamp/requests  ·  자재 /admin/materials  ·  협력사 /admin/partners",
            "계정 /admin/users  ·  서버 /admin/server",
            "",
            "사용 중 오류·권한·데이터 복구는 시스템관리자에게 문의하세요.",
        ],
        base_size=13,
    )

    # ── 마무리 ──
    _title_slide(
        prs,
        "감사합니다",
        "Smart FMS 사용 중 문의는 시스템관리자에게 연락하세요.\n본 매뉴얼은 화면·권한 설정에 따라 일부 메뉴가 다르게 보일 수 있습니다.",
    )

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    return OUT


if __name__ == "__main__":
    path = build()
    print(f"생성 완료: {path} ({path.stat().st_size // 1024} KB)")
    from pptx import Presentation as _P

    n = len(_P(str(path)).slides)
    print(f"슬라이드 수: {n}")
