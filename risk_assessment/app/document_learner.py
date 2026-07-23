"""작업표준서·위험성평가서 문서(Word/PPT/Excel/PDF)에서 학습 데이터 추출"""

from __future__ import annotations

import re
import shutil
import subprocess
import sys
import tempfile
import time
import zipfile
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any

_OLE_MAGIC = b"\xd0\xcf\x11\xe0\xa1\xb1\x1a\xe1"

from app.law_lookup import LAW_DEFAULT

FIVE_M_KEYS = ("Man", "Machine", "Material", "Method", "Management", "Environment")
FIVE_M_LABELS = {
    "man": "Man",
    "machine": "Machine",
    "material": "Material",
    "method": "Method",
    "management": "Management",
    "environment": "Environment",
    "인": "Man",
    "기계": "Machine",
    "자재": "Material",
    "방법": "Method",
    "관리": "Management",
    "환경": "Environment",
}

PHASE_HINTS = {
    "작업 전": ("작업 준비", "작업/준비", "작업준비", "준비", "작업전"),
    "작업 후": (
        "마무리", "정리정돈", "정리 정돈", "현장 정리", "시운전",
        "충수", "마무리작업", "작업후",
    ),
}

META_LABELS = {
    "작업명", "작업구역", "필요한 보호구", "필요한 안전장비", "필요한 자료",
    "필요 공구/장비", "JSA 관리번호", "작업표준 번호", "개정 번호",
    "작성자", "검토자", "승인자", "작성일", "검토일", "승인일",
}

SUPPORTED_SUFFIXES = {".doc", ".docx", ".pptx", ".xlsx", ".xlsm", ".pdf"}

# Word SaveAs2 — XML 문서(.docx), 버전별로 값이 다름
_WD_FORMAT_DOCX_CANDIDATES = (12, 16, 11)


def _is_ole_word_binary(path: Path) -> bool:
    try:
        with path.open("rb") as f:
            return f.read(8) == _OLE_MAGIC
    except OSError:
        return False


def _is_valid_docx_package(path: Path) -> bool:
    if not path.is_file() or path.stat().st_size < 200:
        return False
    if not zipfile.is_zipfile(path):
        return False
    try:
        with zipfile.ZipFile(path) as zf:
            return "[Content_Types].xml" in zf.namelist()
    except (zipfile.BadZipFile, OSError):
        return False


def _find_valid_docx_in_dir(directory: Path) -> Path | None:
    for candidate in sorted(directory.glob("*.docx"), key=lambda p: p.stat().st_mtime, reverse=True):
        if _is_valid_docx_package(candidate):
            return candidate
    return None


def _find_libreoffice_soffice() -> Path | None:
    for candidate in (
        Path(r"C:\Program Files\LibreOffice\program\soffice.exe"),
        Path(r"C:\Program Files (x86)\LibreOffice\program\soffice.exe"),
    ):
        if candidate.is_file():
            return candidate
    found = shutil.which("soffice")
    return Path(found) if found else None


def _convert_doc_with_word(doc_path: Path, out_docx: Path) -> Path:
    import pythoncom
    import win32com.client

    src = str(doc_path.resolve())
    dst = str(out_docx.resolve())
    out_docx.parent.mkdir(parents=True, exist_ok=True)

    pythoncom.CoInitialize()
    word = None
    doc = None
    last_err: Exception | None = None
    try:
        word = win32com.client.DispatchEx("Word.Application")
        word.Visible = False
        word.DisplayAlerts = 0
        doc = word.Documents.Open(
            src,
            ConfirmConversions=False,
            ReadOnly=False,
            AddToRecentFiles=False,
            NoEncodingDialog=True,
        )
        for fmt in _WD_FORMAT_DOCX_CANDIDATES:
            try:
                if out_docx.exists():
                    out_docx.unlink()
                doc.SaveAs2(dst, FileFormat=fmt)
                for _ in range(24):
                    if _is_valid_docx_package(out_docx):
                        return out_docx
                    time.sleep(0.25)
            except Exception as e:
                last_err = e
        alt = _find_valid_docx_in_dir(out_docx.parent)
        if alt:
            if alt != out_docx:
                shutil.copy2(alt, out_docx)
            return out_docx
        raise RuntimeError(last_err or "Word가 docx 파일을 만들지 못했습니다.")
    finally:
        if doc is not None:
            try:
                doc.Close(False)
            except Exception:
                pass
        if word is not None:
            try:
                word.Quit(SaveChanges=0)
            except Exception:
                pass
        pythoncom.CoUninitialize()


def _convert_doc_with_libreoffice(doc_path: Path, out_dir: Path) -> Path:
    soffice = _find_libreoffice_soffice()
    if not soffice:
        raise RuntimeError("LibreOffice(soffice)를 찾을 수 없습니다.")
    out_dir.mkdir(parents=True, exist_ok=True)
    cmd = [
        str(soffice),
        "--headless",
        "--convert-to",
        "docx",
        "--outdir",
        str(out_dir),
        str(doc_path.resolve()),
    ]
    proc = subprocess.run(cmd, capture_output=True, text=True, timeout=120)
    if proc.returncode != 0:
        err = (proc.stderr or proc.stdout or "").strip()
        raise RuntimeError(f"LibreOffice 변환 실패: {err or proc.returncode}")
    out_docx = out_dir / f"{doc_path.stem}.docx"
    if not out_docx.is_file():
        alt = _find_valid_docx_in_dir(out_dir)
        if alt:
            return alt
        raise RuntimeError("LibreOffice 변환 결과(.docx) 파일이 없습니다.")
    if not _is_valid_docx_package(out_docx):
        alt = _find_valid_docx_in_dir(out_dir)
        if alt and _is_valid_docx_package(alt):
            return alt
        raise RuntimeError("LibreOffice 변환 결과가 올바른 docx 형식이 아닙니다.")
    return out_docx


def _convert_pdf_with_libreoffice(pdf_path: Path, out_dir: Path) -> Path:
    """PDF → docx (LibreOffice). 표 추출이 어려운 PDF 보조."""
    soffice = _find_libreoffice_soffice()
    if not soffice:
        raise RuntimeError("LibreOffice(soffice)를 찾을 수 없습니다.")
    out_dir.mkdir(parents=True, exist_ok=True)
    cmd = [
        str(soffice),
        "--headless",
        "--convert-to",
        "docx",
        "--outdir",
        str(out_dir),
        str(pdf_path.resolve()),
    ]
    proc = subprocess.run(cmd, capture_output=True, text=True, timeout=180)
    if proc.returncode != 0:
        err = (proc.stderr or proc.stdout or "").strip()
        raise RuntimeError(f"LibreOffice PDF 변환 실패: {err or proc.returncode}")
    out_docx = out_dir / f"{pdf_path.stem}.docx"
    if not out_docx.is_file():
        alt = _find_valid_docx_in_dir(out_dir)
        if alt:
            return alt
        raise RuntimeError("LibreOffice PDF→docx 결과 파일이 없습니다.")
    return out_docx


def _guess_job_name(
    text_lines: list[str],
    table_rows: list[list[str]],
    default: str,
) -> str:
    job_name = default
    for line in text_lines:
        if re.search(r"작\s*업\s*명", line):
            m = re.split(r"[:：]", line, maxsplit=1)
            if len(m) > 1 and _norm(m[1]):
                return _norm(m[1])
            m2 = re.search(r"작\s*업\s*명\s*[:：]?\s*(.+)", line)
            if m2:
                return _norm(m2.group(1))
    for cells in table_rows:
        joined = " ".join(cells)
        if "작업명" in joined or "작 업 명" in joined:
            for i, c in enumerate(cells):
                if "작업명" in c.replace(" ", "") and i + 1 < len(cells):
                    v = _norm(cells[i + 1])
                    if v:
                        return v
            m = re.search(r"작\s*업\s*명\s*[:：]?\s*(.+)", joined)
            if m:
                return _norm(m.group(1))
    return job_name


def _enrich_five_m_for_learn(
    reg_name: str,
    five_m: dict[str, str],
    rows: list[dict],
    major_name: str | None = None,
) -> tuple[dict[str, str], str]:
    try:
        from app.work_type_lookup import WorkTypeLookup

        return WorkTypeLookup().enrich_five_m_one_e_from_similar(
            reg_name, major_name, five_m, rows
        )
    except Exception:
        return five_m, ""


def _build_word_like_result(
    path: Path,
    *,
    all_table_rows: list[list[str]],
    text_lines: list[str],
    extra_warnings: list[str] | None = None,
    major_name: str | None = None,
) -> list[DocumentLearnResult]:
    reg_name = registration_name_from_path(path)
    merged_rows = _merge_fragmented_table_rows(all_table_rows)
    rows, _ = _rows_from_jsa_table(merged_rows, job_name=reg_name)
    extracted = _extract_five_m_from_table(merged_rows) or _extract_five_m_from_lines(text_lines)
    five_m, infer_note = _enrich_five_m_for_learn(reg_name, extracted, rows, major_name)
    warnings: list[str] = list(extra_warnings or [])
    if infer_note:
        warnings.append(f"5M1E: {infer_note}")
    if not rows:
        warnings.append("문서에서 위험성평가 표 행을 찾지 못했습니다. 5M1E만 저장됩니다.")
    return [
        DocumentLearnResult(
            source_path=str(path),
            job_name=reg_name,
            five_m_one_e=five_m,
            rows=rows,
            warnings=warnings,
        )
    ]


def _collapse_file_results(
    path: Path,
    results: list[DocumentLearnResult],
    major_name: str | None = None,
) -> list[DocumentLearnResult]:
    """한 파일 → 파일명 소분류 1건, 모든 단계·유해·대책 행 병합."""
    reg_name = registration_name_from_path(path)
    all_rows: list[dict] = []
    five_m: dict[str, str] = {}
    warnings: list[str] = []
    for r in results:
        all_rows.extend(r.rows)
        five_m.update(r.five_m_one_e)
        warnings.extend(r.warnings)
    rows = _post_process_learned_rows(all_rows)
    five_m, infer_note = _enrich_five_m_for_learn(reg_name, five_m, rows, major_name)
    if infer_note and f"5M1E: {infer_note}" not in warnings:
        warnings.append(f"5M1E: {infer_note}")
    return [
        DocumentLearnResult(
            source_path=str(path.resolve()),
            job_name=reg_name,
            five_m_one_e=five_m,
            rows=rows,
            warnings=list(dict.fromkeys(warnings)),
        )
    ]


def _needs_legacy_conversion(path: Path) -> bool:
    """구형 Word 바이너리(.doc 또는 docx 확장자의 OLE 파일)."""
    suf = path.suffix.lower()
    if suf == ".doc":
        return True
    if suf == ".docx" and _is_ole_word_binary(path):
        return True
    return False


def _convert_doc_to_docx(doc_path: Path) -> tuple[Path, Path | None]:
    """레거시 .doc → 임시 .docx. 반환: (docx경로, 삭제할 임시폴더 또는 None)."""
    tmp_dir = Path(tempfile.mkdtemp(prefix="pwide_doc_"))
    out_docx = tmp_dir / f"{doc_path.stem}.docx"
    errors: list[str] = []

    if sys.platform == "win32":
        try:
            _convert_doc_with_word(doc_path, out_docx)
            if _is_valid_docx_package(out_docx):
                return out_docx, tmp_dir
            errors.append("Microsoft Word: 변환 파일이 docx(zip) 형식이 아닙니다.")
        except ImportError:
            errors.append("pywin32 미설치 (pip install pywin32)")
        except Exception as e:
            errors.append(f"Microsoft Word: {e}")

    try:
        converted = _convert_doc_with_libreoffice(doc_path, tmp_dir)
        if _is_valid_docx_package(converted):
            return converted, tmp_dir
        errors.append("LibreOffice: 변환 파일이 docx(zip) 형식이 아닙니다.")
    except Exception as e:
        errors.append(f"LibreOffice: {e}")

    shutil.rmtree(tmp_dir, ignore_errors=True)
    hint = (
        "구형 Word(.doc) 파일은 PC에 Microsoft Word 또는 LibreOffice가 필요합니다.\n"
        "· Word 설치: pip install pywin32 후 다시 시도\n"
        "· LibreOffice 설치: https://www.libreoffice.org\n"
        "· 또는 Word에서 「다른 이름으로 저장」→ .docx 후 학습"
    )
    raise RuntimeError(f"『{doc_path.name}』 변환 실패\n" + "\n".join(errors) + f"\n\n{hint}")


def _prepare_docx_path(path: Path) -> tuple[Path, Path | None]:
    """python-docx가 읽을 수 있는 docx 경로로 준비."""
    if _needs_legacy_conversion(path):
        return _convert_doc_to_docx(path)
    if not _is_valid_docx_package(path):
        raise RuntimeError(
            f"『{path.name}』은(는) 올바른 Word(.docx) 파일이 아닙니다.\n"
            "구형 .doc 파일이면 확장자가 .doc인지 확인하거나 .docx로 저장 후 다시 시도하세요."
        )
    return path, None


def _cleanup_temp_dir(temp_dir: Path | None) -> None:
    if temp_dir and temp_dir.is_dir():
        shutil.rmtree(temp_dir, ignore_errors=True)


@dataclass
class DocumentLearnResult:
    source_path: str
    job_name: str
    five_m_one_e: dict[str, str] = field(default_factory=dict)
    rows: list[dict[str, Any]] = field(default_factory=list)
    warnings: list[str] = field(default_factory=list)
    sheet_title: str = ""  # Excel 시트명 — 소분류 등록명으로 사용

    @property
    def row_count(self) -> int:
        return len(self.rows)


def _norm(s: str) -> str:
    if s is None:
        return ""
    s = re.sub(r"\s+", " ", str(s).replace("\n", " ").replace("\u00a0", " ")).strip()
    return re.sub(r"^[∙•·]\s*", "", s)


def _norm_cell(s: str) -> str:
    """표 셀 — 줄바꿈·번호 목록 유지 (유해위험/안전대책 전체 보존)."""
    if s is None:
        return ""
    text = str(s).replace("\r\n", "\n").replace("\r", "\n").replace("\u00a0", " ")
    lines: list[str] = []
    for raw in text.split("\n"):
        line = re.sub(r"\s+", " ", raw.strip())
        line = re.sub(r"^[∙•·]\s*", "", line)
        if line:
            lines.append(line)
    return "\n".join(lines)


def _phase_from_step(step: str) -> str:
    step_n = step.replace(" ", "").replace("/", "")
    for phase, hints in PHASE_HINTS.items():
        for h in hints:
            if h.replace(" ", "") in step_n or h in step:
                return phase
    return "작업 중"


def _parse_int(val: Any, default: int = 0) -> int:
    if val is None:
        return default
    if isinstance(val, (int, float)):
        try:
            n = int(val)
            return n if n > 0 else default
        except (TypeError, ValueError):
            return default
    s = _norm(str(val))
    if not s:
        return default
    m = re.search(r"\b([1-5])\b", s)
    if m:
        return int(m.group(1))
    try:
        n = int(float(s))
        return n if n > 0 else default
    except (TypeError, ValueError):
        return default


def registration_name_from_path(path: str | Path) -> str:
    """소분류 등록명 = 파일명(확장자 제외)."""
    return Path(path).stem.strip() or "학습작업"


def sanitize_registration_name(name: str) -> str:
    """소분류 등록에 쓸 이름 정리."""
    s = re.sub(r'[\\/*?:\[\]"<>|]', "_", (name or "").strip())
    return s[:80] or "시트"


def registration_name_for_result(result: DocumentLearnResult) -> str:
    """문서 학습 결과 → 소분류 등록명 (Excel은 시트명 우선)."""
    if (result.sheet_title or "").strip():
        return sanitize_registration_name(result.sheet_title)
    return registration_name_from_path(result.source_path)


def _row_dict(
    *,
    work_class: str = "",
    phase: str = "",
    unit_task: str = "",
    hazard: str = "",
    injury: str = "",
    current: str = "",
    improvements: str = "",
    freq_before: int = 0,
    sev_before: int = 0,
    freq_after: int = 0,
    sev_after: int = 0,
    law: str = "",
    law_url: str = "",
) -> dict[str, Any]:
    step = _norm(unit_task) or "작업"
    ph = phase or _phase_from_step(step)
    fb = freq_before or 3
    sb = sev_before or 2
    fa = freq_after or max(1, fb - 1)
    sa = sev_after or max(1, sb - 1)
    imp = _norm(improvements)
    cur = _norm(current) or imp
    return {
        "work_class": work_class or "운전작업",
        "phase": ph,
        "unit_task": step,
        "hazard": _norm(hazard),
        "injury": injury or "",
        "current": cur,
        "freq_before": fb,
        "sev_before": sb,
        "improvements": imp or cur,
        "law": law or LAW_DEFAULT[0],
        "law_url": law_url or LAW_DEFAULT[1],
        "freq_after": fa,
        "sev_after": sa,
        "source": "DOC-LEARN",
    }


def _extract_five_m_from_lines(lines: list[str]) -> dict[str, str]:
    out: dict[str, str] = {}
    for line in lines:
        line = _norm(line)
        if not line:
            continue
        for key, label in (
            ("Man", r"(?:Man|인력|작업자|인)\s*[:：]"),
            ("Machine", r"(?:Machine|기계|설비|장비)\s*[:：]"),
            ("Material", r"(?:Material|자재|물질)\s*[:：]"),
            ("Method", r"(?:Method|방법|절차|작업방법)\s*[:：]"),
            ("Management", r"(?:Management|관리|안전관리)\s*[:：]"),
            ("Environment", r"(?:Environment|환경|작업환경)\s*[:：]"),
        ):
            m = re.search(label + r"\s*(.+)$", line, re.I)
            if m and key not in out:
                out[key] = _norm(m.group(1))
    return out


def _extract_five_m_from_table(rows: list[list[str]]) -> dict[str, str]:
    out: dict[str, str] = {}
    jsa_meta = {
        "작업구역": "Environment",
        "필요한보호구": "Man",
        "필요한안전장비": "Machine",
        "필요공구": "Machine",
        "필요한자료": "Method",
        "필요공구장비": "Machine",
    }
    for cells in rows:
        if len(cells) < 2:
            continue
        joined = " ".join(_norm(c) for c in cells[:3])
        label = _norm(cells[0]).replace(" ", "")
        val = _norm_cell(cells[1]) if len(cells) > 1 else ""
        if not val and len(cells) > 2:
            val = _norm_cell(cells[2])
        for meta_key, field in jsa_meta.items():
            if meta_key in label.replace(" ", "") or meta_key in joined.replace(" ", ""):
                if field == "Machine" and out.get("Machine"):
                    out[field] = out["Machine"] + " / " + val
                else:
                    out[field] = val
        for pat, key in FIVE_M_LABELS.items():
            if pat in label.lower() or pat in label:
                out[key] = val
                break
    if out.get("Man") and out.get("Machine"):
        out["Management"] = out.get("Management") or "작업표준·JSA·안전작업허가 준수"
    return out


def _is_header_row(cells: list[str]) -> bool:
    joined = " ".join(cells)
    return ("유해" in joined and "위험" in joined) or (
        "작업단계" in joined and ("유해" in joined or "Hazard" in joined)
    ) or (
        "작업순서" in joined and "유해" in joined
    ) or (
        "안전작업대책" in joined.replace(" ", "")
    )


def _parse_header_map(cells: list[str]) -> dict[str, int]:
    """표 헤더 → 작업단계·유해위험·안전대책 열 인덱스."""
    m: dict[str, int] = {}
    for i, raw in enumerate(cells):
        c = _norm(raw)
        cn = c.replace(" ", "")
        if not cn:
            continue
        if any(k in cn for k in ("작업단계", "작업순서", "Steps")) and "유해" not in cn:
            m.setdefault("step", i)
        elif "작업공정" in cn or cn == "공정":
            m["process"] = i
        elif ("유해" in cn and "위험" in cn) or "유해위험요인" in cn:
            m["hazard"] = i
        elif "재해" in cn:
            m["injury"] = i
        elif "안전작업대책" in cn or ("안전" in cn and "대책" in cn):
            m["safety"] = i
        elif "안전대책" in cn:
            m["safety"] = i
        elif "조치확인" in cn or cn == "확인":
            m["confirm"] = i
        elif "위험도" in cn and "유해" not in cn:
            m["risk_level"] = i
        elif cn == "번호" or cn == "no":
            m["no"] = i
        elif "개선" in cn and ("대책" in cn or "조치" in cn):
            m["after_meas"] = i
        elif "현재" in cn and "안전" in cn:
            m["before_meas"] = i
        elif "안전조치" in cn or "대책" in cn:
            if "before_meas" not in m:
                m["before_meas"] = i
            elif "after_meas" not in m:
                m["after_meas"] = i
            elif "safety" not in m:
                m["safety"] = i
        elif cn in ("F", "빈도") and "before_f" not in m:
            m["before_f"] = i
        elif cn in ("S", "강도") and "before_s" not in m:
            m["before_s"] = i
        elif "법적" in cn or "법령" in cn:
            m["law"] = i
    return m


def _cell_at(cells: list[str], idx: int | None) -> str:
    if idx is None or idx < 0 or idx >= len(cells):
        return ""
    return _norm(cells[idx])


def _cell_at_raw(cells: list[str], idx: int | None) -> str:
    if idx is None or idx < 0 or idx >= len(cells):
        return ""
    return _norm_cell(cells[idx])


def _infer_jsa_indices(header_map: dict[str, int], ncol: int) -> dict[str, int]:
    """번호·작업단계·유해·안전작업대책·빈도·강도 열 추정."""
    m = dict(header_map)
    if "hazard" in m:
        return m
    if ncol >= 4:
        m.setdefault("no", 0)
        m.setdefault("step", 1)
        m.setdefault("hazard", 2)
        m.setdefault("safety", 3)
        if ncol >= 6:
            m.setdefault("before_f", 4)
            m.setdefault("before_s", 5)
    return m


def _merge_fragmented_table_rows(table_rows: list[list[str]]) -> list[list[str]]:
    """PDF에서 쪼개진 표 행을 이전 작업단계 행에 병합."""
    merged: list[list[str]] = []
    for cells in table_rows:
        raw = [str(c) if c is not None else "" for c in cells]
        if not any(x.strip() for x in raw):
            continue
        no = _norm(raw[0]) if raw else ""
        is_numbered = bool(re.match(r"^\d{1,2}$", no.replace(".", "")))
        step_hint = _norm(raw[1]) if len(raw) > 1 else ""
        if not is_numbered and merged and step_hint and "유해" not in step_hint:
            prev = merged[-1]
            for i in range(max(len(prev), len(raw))):
                while len(prev) <= i:
                    prev.append("")
                if i < len(raw) and raw[i].strip():
                    prev[i] = _merge_text_blocks(prev[i], raw[i])
            continue
        if not is_numbered and merged and not step_hint:
            prev = merged[-1]
            for i in range(2, max(len(prev), len(raw))):
                while len(prev) <= i:
                    prev.append("")
                if i < len(raw) and raw[i].strip():
                    prev[i] = _merge_text_blocks(prev[i], raw[i])
            continue
        merged.append(raw)
    return merged


def _merge_text_blocks(a: str, b: str) -> str:
    a, b = (a or "").strip(), (b or "").strip()
    if not a:
        return b
    if not b or b in a:
        return a
    if a in b:
        return b
    return a + "\n" + b


def _post_process_learned_rows(rows: list[dict]) -> list[dict]:
    """작업단계별 1행 유지 — 같은 단계는 유해·대책 본문 병합(누락 방지)."""
    by_step: dict[str, dict] = {}
    order: list[str] = []
    for r in rows:
        step = _norm(r.get("unit_task", "")) or "작업"
        hazard = _norm_cell(r.get("hazard", ""))
        if not hazard or len(hazard.replace(" ", "")) < 2:
            continue
        if any(k in hazard for k in ("MNT-", "JSA 관리번호", "작성일", "Hazard")):
            continue
        imp = _norm_cell(r.get("improvements", ""))
        cur = _norm_cell(r.get("current", ""))
        if step not in by_step:
            by_step[step] = dict(r)
            by_step[step]["unit_task"] = step
            by_step[step]["hazard"] = hazard
            by_step[step]["improvements"] = imp
            by_step[step]["current"] = cur or imp
            order.append(step)
        else:
            prev = by_step[step]
            prev["hazard"] = _merge_text_blocks(prev.get("hazard", ""), hazard)
            prev["improvements"] = _merge_text_blocks(prev.get("improvements", ""), imp)
            prev["current"] = _merge_text_blocks(prev.get("current", ""), cur)
            for nk in ("freq_before", "sev_before", "freq_after", "sev_after"):
                if not prev.get(nk) and r.get(nk):
                    prev[nk] = r[nk]
    out: list[dict] = []
    for step in order:
        r = by_step[step]
        if not r.get("improvements") and r.get("current"):
            r["improvements"] = r["current"]
        if not r.get("current") and r.get("improvements"):
            r["current"] = r["improvements"]
        out.append(r)
    return out


def _is_subheader(cells: list[str]) -> bool:
    if any("유해위험요인" in c and "Hazards" in c for c in cells[:4]):
        return True
    if len(cells) > 2 and not cells[1] and not cells[2]:
        if any(c in ("빈도", "강도", "위험도", "F", "S", "R") for c in cells):
            return True
    return False


def _rows_from_jsa_table(
    table_rows: list[list[str]],
    *,
    job_name: str,
    default_step: str = "",
) -> tuple[list[dict], str]:
    """작업단계·유해위험요인·안전대책 열을 읽어 단계별 행 생성."""
    rows: list[dict] = []
    current_step = default_step
    in_table = False
    header_map: dict[str, int] = {}

    for cells in table_rows:
        cells = [_norm(c) for c in cells]
        if not any(cells):
            continue

        joined = " ".join(cells)
        if "작업명" in joined or "작 업 명" in joined:
            for i, c in enumerate(cells):
                if "작업명" in c.replace(" ", "") and i + 1 < len(cells):
                    found = _norm(cells[i + 1])
                    if found:
                        job_name = found

        if _is_header_row(cells):
            in_table = True
            header_map = _parse_header_map(cells)
            continue

        if not in_table:
            continue
        if _is_subheader(cells):
            continue

        hm = _infer_jsa_indices(header_map, len(cells)) if header_map else {}
        if hm and "hazard" in hm:
            hazard = _cell_at_raw(cells, hm.get("hazard"))
            if len(hazard.replace(" ", "")) < 2 or hazard in META_LABELS:
                pass
            elif any(k in hazard for k in ("MNT-", "JSA 관리번호", "작성일", "Hazard")):
                pass
            else:
                step = _cell_at(cells, hm.get("step")) or _cell_at(cells, hm.get("process"))
                if step and step not in META_LABELS and len(step) < 120:
                    current_step = step
                unit_step = current_step or step or job_name
                injury = _cell_at(cells, hm.get("injury"))
                before_meas = _cell_at_raw(cells, hm.get("before_meas"))
                after_meas = (
                    _cell_at_raw(cells, hm.get("after_meas"))
                    or _cell_at_raw(cells, hm.get("safety"))
                )
                safety = after_meas or before_meas
                fb = _parse_int(_cell_at(cells, hm.get("before_f")))
                sb = _parse_int(_cell_at(cells, hm.get("before_s")))
                law = _cell_at(cells, hm.get("law"))
                rows.append(
                    _row_dict(
                        unit_task=unit_step,
                        phase=_phase_from_step(unit_step),
                        hazard=hazard,
                        injury=injury,
                        current=before_meas or safety,
                        improvements=safety or before_meas,
                        freq_before=fb,
                        sev_before=sb,
                        law=law,
                    )
                )
                continue

        if in_table and len(cells) >= 4:
            no = _norm(cells[0])
            if re.match(r"^\d{1,2}$", no.replace(".", "")):
                step = _norm_cell(cells[1])
                hazard = _norm_cell(cells[2])
                ctrl = _norm_cell(cells[3])
                if step and hazard and "유해" not in step:
                    fb = _parse_int(cells[4]) if len(cells) > 4 else 0
                    sb = _parse_int(cells[5]) if len(cells) > 5 else 0
                    rows.append(
                        _row_dict(
                            unit_task=step,
                            phase=_phase_from_step(step),
                            hazard=hazard,
                            current=ctrl,
                            improvements=ctrl,
                            freq_before=fb,
                            sev_before=sb,
                        )
                    )
                    current_step = step
                    continue

        # PPT/JSA 3열: 작업단계 | 유해위험요인 | 안전대책
        col0 = _norm(cells[0]) if cells else ""
        if col0 and col0 not in META_LABELS and "Steps" not in col0 and "작업" not in col0[:2]:
            if len(col0) >= 2 and not re.fullmatch(r"[\d\-./]+", col0):
                current_step = col0.replace(" / ", "·")

        if len(cells) >= 3:
            a, b, c = _norm(cells[0]), _norm_cell(cells[1]), _norm_cell(cells[2])
            if len(b.replace(" ", "")) >= 2 and b not in META_LABELS and "유해" not in b:
                step = current_step or job_name
                if (
                    a
                    and len(a) >= 2
                    and a not in META_LABELS
                    and a != b
                    and "유해" not in a
                    and not re.fullmatch(r"[\d\-./]+", a)
                ):
                    current_step = a.replace(" / ", "·")
                    step = current_step
                rows.append(
                    _row_dict(
                        unit_task=step,
                        phase=_phase_from_step(step),
                        hazard=b,
                        current=c,
                        improvements=c,
                    )
                )

    return _post_process_learned_rows(rows), job_name


def _split_controls(text: str) -> str:
    text = _norm(text)
    text = text.replace(" ∙ ", "\n").replace("∙", "\n")
    parts = re.split(r"\s*/\s*", text)
    lines = [_norm(p) for p in parts if _norm(p)]
    return "\n".join(lines)


def _parse_pptx(path: Path, major_name: str | None = None) -> list[DocumentLearnResult]:
    from pptx import Presentation

    prs = Presentation(str(path))
    by_job: dict[str, list[dict]] = {}
    five_m_by_job: dict[str, dict[str, str]] = {}
    current_job = path.stem
    current_step = ""
    in_table = False

    for slide in prs.slides:
        slide_lines: list[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                slide_lines.extend(shape.text.splitlines())
            if not getattr(shape, "has_table", False):
                continue
            table = shape.table
            for r in range(len(table.rows)):
                cells = [_norm(table.cell(r, c).text) for c in range(len(table.columns))]
                if not any(cells):
                    continue
                if cells[0] == "작업명" and len(cells) > 1:
                    current_job = _norm(cells[1])
                    by_job.setdefault(current_job, [])
                    current_step = ""
                    in_table = False
                    continue
                if _is_header_row(cells):
                    in_table = True
                    continue
                if not in_table:
                    continue
                if _is_subheader(cells):
                    continue
                if cells[0] and cells[0] not in META_LABELS and "Steps" not in cells[0]:
                    step = cells[0].replace(" / ", "·")
                    if step not in META_LABELS:
                        current_step = step
                if len(cells) < 3:
                    continue
                hazard, control = _norm(cells[1]), _split_controls(cells[2])
                if len(hazard) < 3 or hazard in META_LABELS:
                    continue
                step = current_step or current_job
                by_job.setdefault(current_job, []).append(
                    _row_dict(
                        unit_task=step,
                        phase=_phase_from_step(step),
                        hazard=hazard,
                        current=control,
                        improvements=control,
                    )
                )
            fm = _extract_five_m_from_table(
                [
                    [_norm(table.cell(rr, cc).text) for cc in range(len(table.columns))]
                    for rr in range(len(table.rows))
                ]
            )
            if fm:
                five_m_by_job.setdefault(current_job, {}).update(fm)
        fm_text = _extract_five_m_from_lines(slide_lines)
        if fm_text:
            five_m_by_job.setdefault(current_job, {}).update(fm_text)

    results: list[DocumentLearnResult] = []
    for job, rows in by_job.items():
        if rows:
            results.append(
                DocumentLearnResult(
                    source_path=str(path),
                    job_name=job,
                    five_m_one_e=five_m_by_job.get(job, {}),
                    rows=rows,
                )
            )
    if not results:
        results.append(
            DocumentLearnResult(
                source_path=str(path),
                job_name=registration_name_from_path(path),
                five_m_one_e=five_m_by_job.get(path.stem, {}),
                warnings=["PPT에서 위험성평가 표를 찾지 못했습니다."],
            )
        )
    return _collapse_file_results(path, results, major_name)


def _parse_docx(path: Path, major_name: str | None = None) -> list[DocumentLearnResult]:
    from docx import Document
    from docx.opc.exceptions import PackageNotFoundError

    if path.suffix.lower() == ".doc" or _is_ole_word_binary(path):
        raise RuntimeError(
            f"『{path.name}』은(는) 구형 Word(.doc) 형식입니다. "
            "변환 단계가 누락되었습니다. 프로그램을 최신 버전으로 실행해 주세요."
        )
    if not _is_valid_docx_package(path):
        raise RuntimeError(
            f"『{path.name}』을(를) Word 문서로 열 수 없습니다.\n"
            "파일이 손상되었거나 .doc(구형) 파일일 수 있습니다.\n"
            "Microsoft Word 또는 LibreOffice에서 .docx로 저장한 뒤 다시 학습해 주세요."
        )
    try:
        doc = Document(str(path.resolve()))
    except PackageNotFoundError as e:
        raise RuntimeError(
            f"『{path.name}』 — Word 패키지를 읽을 수 없습니다.\n"
            "구형 .doc 파일이면 Word/LibreOffice 설치 후 다시 시도하거나 .docx로 저장하세요.\n"
            f"(상세: {e})"
        ) from e
    all_table_rows: list[list[str]] = []
    text_lines: list[str] = []

    for p in doc.paragraphs:
        t = _norm(p.text)
        if t:
            text_lines.append(t)

    for table in doc.tables:
        for row in table.rows:
            cells = [_norm(c.text) for c in row.cells]
            all_table_rows.append(cells)

    return _build_word_like_result(
        path, all_table_rows=all_table_rows, text_lines=text_lines, major_name=major_name
    )


def _parse_pdf(path: Path, major_name: str | None = None) -> list[DocumentLearnResult]:
    try:
        import pdfplumber
    except ImportError as e:
        raise ImportError(
            "PDF 학습에 pdfplumber가 필요합니다. pip install pdfplumber"
        ) from e

    all_table_rows: list[list[str]] = []
    text_lines: list[str] = []
    page_count = 0

    table_settings_list = (
        {"vertical_strategy": "lines", "horizontal_strategy": "lines", "intersection_tolerance": 8},
        {"vertical_strategy": "text", "horizontal_strategy": "text"},
    )
    with pdfplumber.open(str(path.resolve())) as pdf:
        if not pdf.pages:
            raise RuntimeError(f"『{path.name}』 PDF에 페이지가 없습니다.")
        page_count = len(pdf.pages)
        for page in pdf.pages:
            text = page.extract_text() or ""
            for line in text.splitlines():
                t = _norm(line)
                if t:
                    text_lines.append(t)
            seen_sig: set[str] = set()
            for settings in table_settings_list:
                try:
                    tables = page.extract_tables(table_settings=settings) or []
                except Exception:
                    tables = page.extract_tables() or []
                for table in tables:
                    for row in table:
                        if not row:
                            continue
                        raw_cells = [str(c) if c is not None else "" for c in row]
                        if not any(c.strip() for c in raw_cells):
                            continue
                        sig = "|".join(_norm(c)[:30] for c in raw_cells[:4])
                        if sig in seen_sig:
                            continue
                        seen_sig.add(sig)
                        all_table_rows.append(raw_cells)

    extra: list[str] = []
    if page_count and len("".join(text_lines)) < 40 and not all_table_rows:
        extra.append(
            "스캔(이미지) PDF로 보입니다. 텍스트·표 추출이 어렵습니다. "
            "Word/Excel로 변환하거나 OCR 후 다시 시도하세요."
        )

    result = _build_word_like_result(
        path,
        all_table_rows=all_table_rows,
        text_lines=text_lines,
        extra_warnings=extra,
        major_name=major_name,
    )

    if result[0].rows:
        return result

    if _find_libreoffice_soffice():
        tmp_dir = Path(tempfile.mkdtemp(prefix="pwide_pdf_"))
        try:
            docx_path = _convert_pdf_with_libreoffice(path, tmp_dir)
            if _is_valid_docx_package(docx_path):
                sub = _parse_docx(docx_path, major_name)
                if sub and sub[0].rows:
                    sub[0].source_path = str(path.resolve())
                    sub[0].warnings = sub[0].warnings or []
                    sub[0].warnings.insert(0, "PDF — LibreOffice docx 변환 후 분석")
                    return sub
        except Exception as e:
            result[0].warnings.append(f"PDF docx 변환 보조 실패: {e}")
        finally:
            _cleanup_temp_dir(tmp_dir)

    return result


def _parse_excel(path: Path, major_name: str | None = None) -> list[DocumentLearnResult]:
    from openpyxl import load_workbook

    wb = load_workbook(str(path), read_only=True, data_only=True)
    results: list[DocumentLearnResult] = []
    used_reg_names: dict[str, int] = {}

    def _unique_reg(base: str) -> str:
        name = sanitize_registration_name(base)
        if name not in used_reg_names:
            used_reg_names[name] = 0
            return name
        used_reg_names[name] += 1
        return f"{name}_{used_reg_names[name]}"

    for sheet in wb.worksheets:
        sheet_title = (sheet.title or "").strip() or "시트"
        reg_name = _unique_reg(sheet_title)
        grid: list[list[str]] = []
        doc_job = ""

        for row in sheet.iter_rows(values_only=True):
            cells = [_norm(v) if v is not None else "" for v in row]
            if not any(cells):
                continue
            grid.append(cells)
            joined = " ".join(cells)
            if ("작업명" in joined or "작 업 명" in joined) and not doc_job:
                for i, c in enumerate(cells):
                    if "작업명" in c.replace(" ", "") and i + 1 < len(cells):
                        doc_job = _norm(cells[i + 1])
                        break
                if not doc_job:
                    m = re.search(r"작\s*업\s*명\s*[:：]?\s*(.+)", joined)
                    if m:
                        doc_job = _norm(m.group(1))

        if not grid:
            continue

        merged = _merge_fragmented_table_rows(grid)
        rows, _ = _rows_from_jsa_table(merged, job_name=reg_name)
        extracted = _extract_five_m_from_table(merged) or _extract_five_m_from_lines(
            [c for row in merged for c in row if c]
        )
        if not rows and not any(extracted.values()):
            continue

        five_m, infer_note = _enrich_five_m_for_learn(reg_name, extracted, rows, major_name)
        warnings: list[str] = [f"Excel 시트: {sheet_title}"]
        if infer_note:
            warnings.append(f"5M1E: {infer_note}")
        if not rows:
            warnings.append("이 시트에서 위험성평가 표 행을 찾지 못았습니다. 5M1E만 저장됩니다.")

        results.append(
            DocumentLearnResult(
                source_path=str(path.resolve()),
                job_name=reg_name,
                sheet_title=sheet_title,
                five_m_one_e=five_m,
                rows=rows,
                warnings=warnings,
            )
        )

    wb.close()
    if not results:
        results.append(
            DocumentLearnResult(
                source_path=str(path.resolve()),
                job_name=registration_name_from_path(path),
                warnings=["Excel에서 학습 가능한 시트를 찾지 못았습니다."],
            )
        )
    return results


def parse_document(path: str | Path, major_name: str | None = None) -> list[DocumentLearnResult]:
    path = Path(path)
    if not path.exists():
        raise FileNotFoundError(str(path))
    suffix = path.suffix.lower()
    if suffix not in SUPPORTED_SUFFIXES:
        raise ValueError(
            f"지원하지 않는 형식: {suffix} (doc, docx, pptx, xlsx, xlsm, pdf)"
        )

    if suffix == ".pdf":
        return _parse_pdf(path, major_name)
    if suffix == ".pptx":
        return _parse_pptx(path, major_name)
    if suffix in (".docx", ".doc"):
        work_path, temp_dir = _prepare_docx_path(path)
        try:
            results = _parse_docx(work_path, major_name)
            if _needs_legacy_conversion(path):
                orig = str(path.resolve())
                for item in results:
                    item.source_path = orig
                    item.warnings = item.warnings or []
                    if not any("변환" in w for w in item.warnings):
                        item.warnings.insert(
                            0, "구형 .doc 파일 — Word/LibreOffice로 변환 후 분석",
                        )
            return results
        finally:
            _cleanup_temp_dir(temp_dir)
    return _parse_excel(path, major_name)


def parse_documents(
    paths: list[str | Path],
    major_name: str | None = None,
) -> list[DocumentLearnResult]:
    merged: dict[str, DocumentLearnResult] = {}
    for p in paths:
        for item in parse_document(p, major_name):
            key = f"{Path(item.source_path).resolve()}::{registration_name_for_result(item)}"
            if key not in merged:
                merged[key] = item
                continue
            prev = merged[key]
            prev.rows = _post_process_learned_rows(prev.rows + item.rows)
            for k, v in item.five_m_one_e.items():
                if v and not prev.five_m_one_e.get(k):
                    prev.five_m_one_e[k] = v
            prev.warnings.extend(item.warnings)
    return list(merged.values())


def results_to_risk_rows(results: list[DocumentLearnResult]):
    from app.local_engine import RiskRow

    out: list[RiskRow] = []
    for res in results:
        for d in res.rows:
            out.append(
                RiskRow(
                    work_class=str(d.get("work_class", "")),
                    phase=str(d.get("phase", "")),
                    unit_task=str(d.get("unit_task", "")),
                    hazard=str(d.get("hazard", "")),
                    injury=str(d.get("injury", "")),
                    current=str(d.get("current", "")),
                    freq_before=int(d.get("freq_before") or 0),
                    sev_before=int(d.get("sev_before") or 0),
                    improvements=str(d.get("improvements", "")),
                    law=str(d.get("law", LAW_DEFAULT[0])),
                    law_url=str(d.get("law_url", "")),
                    freq_after=int(d.get("freq_after") or 0),
                    sev_after=int(d.get("sev_after") or 0),
                    source=str(d.get("source") or "DOC-LEARN"),
                )
            )
    return out
